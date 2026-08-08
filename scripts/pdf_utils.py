import contextlib
import threading
import logging
import re
import uuid
import pikepdf
from pathlib import Path
import os
from typing import List, Optional

from scripts.utils import branded_filename, original_stem

# --- Logging Setup ---
logger = logging.getLogger(__name__)

try:
    import fitz
except ImportError:
    fitz = None
try:
    import cv2
    import numpy as np
except ImportError:
    cv2 = None
    np = None
try:
    from docxcompose.composer import Composer
    from docx import Document as Document_docx
except ImportError:
    Composer = None
    Document_docx = None
import shutil


# Global cache for PaddleOCR engine to avoid expensive re-initialization
_PADDLE_ENGINE = None
_ENGINE_LOCK = threading.Lock()


def _resolve_models_dir() -> Path:
    """Resolve the Paddle models directory from common project layouts."""
    script_dir = Path(__file__).resolve().parent
    candidates = [script_dir / "models", script_dir.parent / "models"]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return candidates[1]

def get_paddle_engine():
    """Returns cached PaddleOCR engine instance, initializing if needed (thread-safe)."""
    global _PADDLE_ENGINE
    if _PADDLE_ENGINE is None:
        with _ENGINE_LOCK:
            if _PADDLE_ENGINE is None:
                logger.info("Initializing PaddleOCR engine (first time only)...")
                # Disable MKL-DNN/OneDNN to fix compatibility issues on Windows.
                # Must be set BEFORE importing paddle/paddleocr, which is why it
                # lives here next to the deferred import (and only runs when the
                # Paddle backend is actually used).
                os.environ['FLAGS_use_mkldnn'] = '0'
                os.environ['MKLDNN_VERBOSE'] = '0'
                os.environ['PADDLE_DISABLE_MKLDNN'] = '1'
                os.environ['FLAGS_enable_mkldnn'] = '0'
                # Force CPU-only mode with basic backend
                os.environ['CUDA_VISIBLE_DEVICES'] = ''
                try:
                    # Deferred imports to prevent boot crashes on low-resource environments
                    from paddleocr import PPStructure
                except ImportError as e:
                    logger.exception("PaddleOCR not installed correctly")
                    raise ImportError("PaddleOCR engine is missing. If you are on the Free Tier, it might have failed to install due to size limits.") from e

                paddle_dir = _resolve_models_dir()
                # When use_onnx=True, PaddleOCR passes *_model_dir directly to
                # ort.InferenceSession as a FILE path (despite the "dir" name).
                # Auto-downloaded ONNX models are single .onnx files; our vendored
                # models are stored as model.onnx inside each subdirectory.
                layout_dir = paddle_dir / "layout" / "picodet_lcnet_x1_0_fgd_layout_infer" / "model.onnx"
                table_dir = paddle_dir / "table" / "en_ppstructure_mobile_v2.0_SLANet_inference" / "model.onnx"
                det_dir = paddle_dir / "det" / "en" / "en_PP-OCRv3_det_infer" / "model.onnx"
                rec_dir = paddle_dir / "rec" / "en" / "en_PP-OCRv3_rec_infer" / "model.onnx"

                try:
                    _PADDLE_ENGINE = PPStructure(
                        recovery=True, lang='en', show_log=False, use_gpu=False,
                        enable_mkldnn=False, use_onnx=True,
                        layout_model_dir=str(layout_dir),
                        table_model_dir=str(table_dir),
                        det_model_dir=str(det_dir),
                        rec_model_dir=str(rec_dir)
                    )
                except MemoryError:
                    logger.critical("Out of Memory while loading PaddleOCR.")
                    raise MemoryError("Server ran out of memory loading the AI engine. Please upgrade to a 'Starter' plan on Render for this feature.")
                except Exception:
                    logger.exception("Unexpected error loading PaddleOCR")
                    raise
                    
                logger.info("PaddleOCR engine cached successfully")
    return _PADDLE_ENGINE

def remove_pdf_password(input_path: str, password: str, output_dir: str) -> str:
    """Remove password protection / restrictions from a PDF and save to output_dir.

    Owner-restricted PDFs carry permission restrictions (no copy/print/edit) but
    no open password, so they open with an *empty* password. When the supplied
    password is rejected we therefore retry once with an empty password before
    giving up: this strips owner-only restrictions even when the caller passes a
    wrong or irrelevant password, matching how mainstream "Unlock PDF" tools
    (iLovePDF, Smallpdf, PDF24) behave. A PDF that genuinely has an *open*
    (user) password still fails on a wrong password, because the empty-password
    retry cannot open it either — so this never produces a false unlock.
    """
    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    try:
        pdf = pikepdf.open(input_file, password=password)
    except pikepdf.PasswordError:
        # Fall back to an empty password: succeeds only for owner-restricted
        # PDFs (no open password). Re-raise the original error otherwise.
        try:
            pdf = pikepdf.open(input_file, password="")
        except pikepdf.PasswordError:
            raise

    with pdf:
        pdf.save(output_file)

    return str(output_file)

def _parse_page_selection(pages: str, total_pages: int) -> List[int]:
    """Parse a page selection string (e.g., '1,3-5') into zero-based indices."""
    if pages is None:
        raise ValueError("No pages selected. Please provide page numbers or 'all'.")

    normalized = pages.strip().lower()
    if not normalized:
        raise ValueError("No pages selected. Please provide page numbers or 'all'.")

    if normalized == "all":
        return list(range(total_pages))

    indices: List[int] = []
    seen = set()

    for part in normalized.split(","):
        segment = part.strip()
        if not segment:
            continue

        if "-" in segment:
            start_str, end_str = segment.split("-", 1)
            if not start_str or not end_str:
                raise ValueError(f"Invalid page range segment: '{segment}'")
            try:
                start = int(start_str)
                end = int(end_str)
            except ValueError:
                raise ValueError(f"Invalid page range numbers: '{segment}'")
            if start < 1 or end < 1 or start > end:
                raise ValueError(f"Invalid page range segment: '{segment}'")
            for num in range(start, end + 1):
                if num not in seen:
                    seen.add(num)
                    indices.append(num - 1)
        else:
            try:
                num = int(segment)
            except ValueError:
                raise ValueError(f"Invalid page number: '{segment}'")
            if num < 1:
                raise ValueError(f"Invalid page number: '{segment}'")
            if num not in seen:
                seen.add(num)
                indices.append(num - 1)

    if not indices:
        raise ValueError("No valid pages selected.")

    if max(indices) >= total_pages:
        raise ValueError(f"Selected page number exceeds document page count ({total_pages}).")

    return indices

def _get_decrypted_pdf_path(input_path: str, password: str = None, temp_dir: Path = None) -> tuple:
    """
    Returns (path_to_use, needs_cleanup).
    If encrypted and password provided, decrypts to temp file.
    If encrypted and no password, raises ValueError.
    If not encrypted, returns original path.
    """
    input_file = Path(input_path)
    
    # Try to open PDF - check if encrypted in a single operation
    try:
        with pikepdf.open(input_file) as pdf:
            # PDF is not encrypted or has no password
            return str(input_file), False
    except pikepdf.PasswordError:
        # PDF is encrypted - need password
        if not password:
            raise ValueError(f"PDF is password-protected. Please provide a password.")
        
        # Decrypt to temp file
        if temp_dir is None:
            temp_dir = input_file.parent
        
        temp_file = temp_dir / f"{original_stem(input_file)}_temp_decrypted.pdf"
        
        with pikepdf.open(input_file, password=password) as pdf:
            pdf.save(temp_file)
        
        return str(temp_file), True

def extract_pdf_pages(input_path: str, output_dir: str, pages: str, password: str = None) -> str:
    """Extract selected pages from PDF and save to output_dir."""
    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        with pikepdf.open(decrypted_path) as pdf:
            selected_indices = _parse_page_selection(pages, len(pdf.pages))

            new_pdf = pikepdf.Pdf.new()
            for idx in selected_indices:
                new_pdf.pages.append(pdf.pages[idx])

            new_pdf.save(output_file)
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return str(output_file)


def _normalize_extracted_text(text: str) -> str:
    """Collapse page text into readable plain text without preserving layout spacing."""
    lines = [" ".join(line.split()) for line in text.splitlines()]
    return "\n".join(line for line in lines if line)


# Tunable thresholds for "does this page already carry usable text".
TEXT_LAYER_MIN_CHARS_PER_PAGE = 40
TEXT_LAYER_MIN_WORDS_PER_PAGE = 6
TEXT_LAYER_MIN_PAGE_FRACTION = 0.8  # doc-level: fraction of pages that must qualify


def _page_has_usable_text(page) -> bool:
    """True if `page`'s embedded text layer (via PyMuPDF) looks like real
    body text rather than stray metadata/watermark noise."""
    words = page.get_text("text").split()
    non_ws_chars = sum(len(w) for w in words)
    return (
        non_ws_chars >= TEXT_LAYER_MIN_CHARS_PER_PAGE
        and len(words) >= TEXT_LAYER_MIN_WORDS_PER_PAGE
    )


def _inspect_text_layer(doc: "fitz.Document") -> dict:
    """Classify each page of an already-open fitz.Document as text/no-text.

    Used to decide, per page, whether rasterize+OCR is actually needed —
    OCR should only run on pages that don't already have a usable
    embedded text layer.
    """
    pages_with_text = [_page_has_usable_text(page) for page in doc]
    total_pages = len(pages_with_text)
    fraction = (sum(pages_with_text) / total_pages) if total_pages else 0.0
    return {
        "pages_with_text": pages_with_text,
        "total_pages": total_pages,
        "fraction_with_text": fraction,
        "has_usable_text_layer": fraction >= TEXT_LAYER_MIN_PAGE_FRACTION,
    }


def _render_page_bgr(page, dpi: int = 200):
    """Render a PDF page to a BGR numpy image for OCR."""
    pix = page.get_pixmap(dpi=dpi)
    img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)

    if pix.n == 3:
        return cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR)
    if pix.n == 4:
        return cv2.cvtColor(img_array, cv2.COLOR_RGBA2BGR)
    return cv2.cvtColor(img_array, cv2.COLOR_GRAY2BGR)


def _extract_text_with_ocr(doc: fitz.Document) -> str:
    """Extract text from rendered PDF pages using the configured OCR backend.

    Returns an empty string when no OCR backend is available (DISABLE_AI=1),
    so callers fall through to their normal "no text" handling.
    """
    from scripts.ocr_engine import get_ocr_engine

    ocr_engine = get_ocr_engine()
    if ocr_engine is None:
        logger.info("OCR fallback skipped: no OCR backend configured")
        return ""

    page_text: List[str] = []
    for page in doc:
        img = _render_page_bgr(page)
        items = ocr_engine.recognize(img)
        page_text.append("\n".join(item["text"] for item in items if item.get("text")))

    return "\n\n".join(text for text in page_text if text.strip())


def extract_pdf_text(
    input_path: str,
    output_dir: str,
    password: str = None,
    preserve_formatting: bool = True,
    use_ocr: bool = True,
    min_text_chars: int = 20,
) -> str:
    """Extract text from a PDF and save it as a .txt file."""
    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "txt")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)
    doc = None

    try:
        doc = fitz.open(decrypted_path)
        page_text: List[str] = []

        for page in doc:
            text = page.get_text("text")
            if not preserve_formatting:
                text = _normalize_extracted_text(text)
            page_text.append(text.strip())

        extracted_text = "\n\n".join(text for text in page_text if text)

        if len(extracted_text.strip()) < min_text_chars and use_ocr:
            extracted_text = _extract_text_with_ocr(doc)

        if not extracted_text.strip():
            raise ValueError("No text could be extracted from this PDF.")

        output_file.write_text(extracted_text.rstrip() + "\n", encoding="utf-8")
        return str(output_file)
    finally:
        if doc is not None:
            doc.close()
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

def compress_pdf(input_path: str, output_dir: str, level: str = 'medium', password: str = None) -> dict:
    """Compress PDF by optimizing structure and resampling large images.
    
    Returns dict with output_path, original_size, compressed_size, reduction_pct.
    """
    import io
    import shutil
    import fitz
    from PIL import Image as PILImage

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    # Per-level image targets: smaller max dimension + lower JPEG quality => smaller
    # file. Every level (including 'low') does real image recompression so the
    # three options produce visibly different results.
    level = (level or 'medium').lower()
    if level not in ('low', 'medium', 'high'):
        level = 'medium'
    max_dim, jpeg_quality = {
        'low':    (2200, 80),
        'medium': (1600, 60),
        'high':   (1000, 40),
    }[level]

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        original_size = input_file.stat().st_size
        doc = fitz.open(decrypted_path)

        xrefs_processed = set()
        for page in doc:
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                if xref in xrefs_processed:
                    continue
                xrefs_processed.add(xref)
                try:
                    # Current on-disk (still-compressed) size of this image stream,
                    # so we only swap in a re-encoded version when it's smaller.
                    try:
                        old_len = len(doc.xref_stream_raw(xref))
                    except Exception:
                        old_len = None

                    pix = fitz.Pixmap(doc, xref)

                    # Not worth re-encoding tiny images (icons, bullets, logos).
                    if pix.width * pix.height < 80 * 80:
                        pix = None
                        continue

                    # Normalise to a plain grayscale/RGB buffer Pillow can read.
                    if pix.alpha:
                        pix = fitz.Pixmap(pix, 0)  # drop the alpha channel
                    if pix.n == 1:
                        mode = "L"
                    elif pix.n == 3:
                        mode = "RGB"
                    else:  # CMYK or other -> convert to RGB
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                        mode = "RGB"

                    pil_img = PILImage.frombytes(mode, (pix.width, pix.height), pix.samples)
                    pix = None
                    if pil_img.mode != 'RGB':
                        pil_img = pil_img.convert('RGB')

                    # Downscale anything larger than this level's max dimension.
                    longest = max(pil_img.width, pil_img.height)
                    if longest > max_dim:
                        factor = max_dim / longest
                        pil_img = pil_img.resize(
                            (max(1, int(pil_img.width * factor)),
                             max(1, int(pil_img.height * factor))),
                            PILImage.LANCZOS,
                        )

                    buf = io.BytesIO()
                    pil_img.save(buf, format='JPEG', quality=jpeg_quality, optimize=True)
                    jpeg_bytes = buf.getvalue()

                    # NOTE: replace_image lives on Page (not Document) in modern
                    # PyMuPDF. Calling it on the document silently failed before,
                    # which is why every level produced identical output.
                    if old_len is None or len(jpeg_bytes) < old_len:
                        page.replace_image(xref, stream=jpeg_bytes)

                except Exception as e:
                    logger.warning("Skipping image xref %s: %s", xref, e)
                    continue

        doc.save(
            str(output_file),
            garbage=4,
            deflate=True,
            deflate_images=True,
            deflate_fonts=True,
            clean=True,
        )
        doc.close()

        compressed_size = output_file.stat().st_size

        # Never hand back a file larger than the input: if the source was already
        # well-optimised, fall back to the (decrypted) original.
        if compressed_size >= original_size:
            shutil.copyfile(decrypted_path, output_file)
            compressed_size = output_file.stat().st_size

        reduction = max(0.0, (1 - compressed_size / original_size) * 100)

        return {
            'output_path': str(output_file),
            'original_size': original_size,
            'compressed_size': compressed_size,
            'reduction_pct': round(reduction, 1),
        }

    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)


# pdf2docx's multi_processing spawns one worker per page batch. Measured on a
# 60-page text PDF (best of 2 runs each):
#
#   4 cores available:  3.7s multiprocessing  vs  9.5s serial  -> pool wins 2.5x
#   1 core available:  10.5s multiprocessing  vs  9.8s serial  -> serial wins ~7%
#
# So the pool is a large win when there are cores to use and a small loss when
# there aren't, because the workers contend for one CPU and each pays the full
# interpreter + PyMuPDF import cost. Production runs on an Always-Free Oracle VM
# (1-2 cores), so it should convert serially.
#
# Note this is a modest effect and does NOT by itself explain the 2622s p95
# /admin/stats recorded — that is far beyond anything reproduced here, and more
# likely a very large input and/or memory pressure on a small VM. Treat this as
# removing a known regression on small boxes, not as the whole fix; re-check the
# p95 on the next snapshot before concluding anything.
_MULTIPROCESSING_MIN_CORES = 4


def _available_cores() -> int:
    """Cores this process may actually run on.

    Not os.cpu_count(): that reports the host's cores, which on a cgroup-limited
    VM or container is the wrong number — measured 4 here while the process was
    pinned to a single CPU. sched_getaffinity is the real budget, and is Linux-
    only, so fall back to cpu_count where it doesn't exist (Windows, macOS)."""
    try:
        return len(os.sched_getaffinity(0))
    except AttributeError:
        return os.cpu_count() or 1


def _use_multiprocessing() -> bool:
    """Whether to hand pdf2docx a worker pool. Env var wins so a deployment can
    override the heuristic without a code change; otherwise decide on cores."""
    override = os.environ.get("PDF2DOCX_MULTIPROCESSING")
    if override is not None and override.strip() != "":
        return override.strip().lower() in ("1", "true", "yes", "on")
    return _available_cores() >= _MULTIPROCESSING_MIN_CORES


# pdf2docx reports its own per-page progress with `logging.info` on the ROOT
# logger, in the fixed form "(3/17) Page 3" — twice per document, once while
# parsing pages and once while writing them. That is the only progress signal
# it exposes: convert() takes no callback. Parsing it back out of a log record
# is not elegant, but the alternative is a 6-minute conversion that streams
# nothing at all, and the SSE endpoint exists precisely to avoid that.
_PDF2DOCX_PAGE_LOG = re.compile(r"^\((\d+)/(\d+)\)\s+Page\s+\d+")


class _Pdf2docxProgress(logging.Handler):
    """Translate pdf2docx's page logs into progress_callback calls.

    Scoped to the calling thread: this attaches to the root logger, so in a
    server handling concurrent conversions it would otherwise see every other
    conversion's pages too and report one request's progress into another's
    stream. pdf2docx's serial path logs from the thread that called it, so
    thread identity is the correct filter.

    How many passes are visible depends on the mode, so the caller says which:

    * serial — both the parsing and the writing pass log in this thread. They
      are folded into one monotonic 0..total count, because reporting each as
      1..total would send the progress bar back to the start halfway through,
      which reads as a failure to anyone watching it.
    * multiprocessing — parsing happens in subprocesses whose logging never
      reaches this process, so only the writing pass is visible and it maps
      straight through.
    """

    def __init__(self, progress_callback, single_pass: bool):
        super().__init__(level=logging.INFO)
        self._callback = progress_callback
        self._single_pass = single_pass
        self._thread = threading.get_ident()
        self._seen_first_pass = False
        self._furthest = 0

    def emit(self, record):
        # A logging handler that raises would take the conversion down with
        # it — this is progress reporting, it must never be load-bearing.
        try:
            if record.thread != self._thread:
                return
            match = _PDF2DOCX_PAGE_LOG.match(record.getMessage())
            if not match:
                return
            done, total = int(match.group(1)), int(match.group(2))
            if self._single_pass:
                _safe_progress(self._callback, min(done, total), total)
                return
            if done < self._furthest:
                self._seen_first_pass = True  # counter restarted: second pass
            self._furthest = done
            half = total / 2
            scaled = (half + done / 2) if self._seen_first_pass else (done / 2)
            _safe_progress(self._callback, min(int(scaled), total), total)
        except Exception:  # noqa: BLE001
            pass


@contextlib.contextmanager
def _pdf2docx_progress(progress_callback, single_pass: bool):
    """Attach the progress handler for the duration of a conversion."""
    if progress_callback is None:
        yield
        return
    handler = _Pdf2docxProgress(progress_callback, single_pass)
    root = logging.getLogger()
    # pdf2docx logs at INFO; if the root logger is above that the records are
    # dropped before any handler sees them, so make sure they get through.
    previous_level = root.level
    if previous_level > logging.INFO:
        root.setLevel(logging.INFO)
    root.addHandler(handler)
    try:
        yield
    finally:
        root.removeHandler(handler)
        if previous_level > logging.INFO:
            root.setLevel(previous_level)


# A missing *system* library surfaces as an ImportError naming a .so file —
# "libGL.so.1: cannot open shared object file". pdf2docx imports cv2, so a host
# that never installed the OpenCV runtime libraries fails here rather than at
# startup, once per conversion, for as long as nobody notices.
_MISSING_SHARED_LIB = re.compile(r"(lib[\w.+-]*\.so[\w.]*)")


class ServerDependencyError(ImportError):
    """A required component is missing from the *server*, not the user's file.

    Raised instead of letting a raw ImportError reach the user. The two read
    very differently: "libGL.so.1: cannot open shared object file" invites
    someone to go and re-save their PDF, which cannot possibly help, while this
    says the problem is ours and their file was never the issue.

    Subclasses ImportError rather than RuntimeError so it stays exactly what it
    replaced: callers and tests that treat a missing dependency as an
    ImportError keep working, and only the message and the HTTP status change.
    """


def _server_dependency_error(feature: str, exc: Exception) -> "ServerDependencyError":
    """Translate a missing-dependency ImportError, and log it for the operator."""
    match = _MISSING_SHARED_LIB.search(str(exc))
    library = match.group(1) if match else None
    logger.error(
        "%s is unavailable: a required dependency is missing on this host (%s). "
        "This is a provisioning gap, not a code bug — see "
        "docs/fix-libreoffice-libgl-oracle.md for the runbook.",
        feature, exc,
    )
    detail = f" (missing system library {library})" if library else ""
    return ServerDependencyError(
        f"{feature} is temporarily unavailable on this server{detail}. "
        "This is a server configuration problem, not a problem with your file — "
        "nothing you change about the document will help. Please try again later."
    )


# Once a multiprocessing pass has failed in this process it will keep failing —
# the causes are environmental (no fork, no /dev/shm, a cgroup CPU budget), not
# per-document. Remembering it turns "every conversion silently costs two full
# conversions" into "the first one does".
_multiprocessing_broken = False


def _convert_pdf2docx(decrypted_path: str, output_file: Path,
                      progress_callback=None) -> None:
    """Shared pdf2docx conversion core, used by both the standard converter
    and pdf_to_word_ai's automatic text-layer routing.

    `progress_callback(page_done, total_pages)` is driven from pdf2docx's own
    page logging. Under multiprocessing only the writing pass is visible (the
    parsing pass runs in subprocesses), which is the right way round: that mode
    is the fast one, and the slow serial path — where a user most needs to be
    told something is still happening — reports both passes.
    """
    global _multiprocessing_broken
    # pdf2docx pulls in cv2, which needs OpenCV's runtime shared libraries. On a
    # host missing them this is where every PDF→Word conversion dies, and the
    # raw message names a .so file the user can do nothing about. Both the
    # import and the construction are covered: which of the two raises depends
    # on whether the dependency underneath is imported eagerly or lazily.
    try:
        from pdf2docx import Converter

        cv = Converter(decrypted_path)
    except ImportError as exc:
        raise _server_dependency_error("PDF to Word conversion", exc) from exc

    multi = _use_multiprocessing() and not _multiprocessing_broken
    try:
        try:
            with _pdf2docx_progress(progress_callback, single_pass=multi):
                cv.convert(str(output_file), multi_processing=multi)
        except ImportError as exc:
            # A lazily-imported dependency can fail here rather than above.
            # Caught before the retry below so it surfaces on the first attempt
            # instead of doubling the wait before the same failure.
            raise _server_dependency_error("PDF to Word conversion", exc) from exc
        except (OSError, RuntimeError, ValueError, AssertionError) as exc:
            # Some environments (e.g. spawn-only Windows workers) can't fork —
            # fall back to a serial pass. Deliberately narrow: this retry costs
            # a second full conversion, so it must not swallow failures that
            # will simply happen again. In particular an ImportError (missing
            # libGL.so.1 behind pdf2docx's cv2 import) has to surface on the
            # first attempt instead of doubling the time before it does.
            if not multi:
                raise
            logger.warning(
                "pdf2docx multiprocessing pass failed (%s: %s) — retrying "
                "serially, and using the serial path directly from now on",
                type(exc).__name__, exc,
            )
            _multiprocessing_broken = True
            with _pdf2docx_progress(progress_callback, single_pass=False):
                cv.convert(str(output_file))
    finally:
        # pdf2docx Converter holds the source PDF open via PyMuPDF until close()
        # is called. Without this finally block, a failed convert() would leak
        # the file handle and block cleanup on Windows.
        cv.close()


def pdf_to_docx(input_path: str, output_dir: str, password: str = None,
                progress_callback=None) -> str:
    """Converts PDF to DOCX using pdf2docx (Fast, Rule-based).

    `progress_callback(page_done, total_pages)` mirrors pdf_to_word_ai's, so
    the streaming endpoint can report progress on this path too. It could not
    before: the standard converter is the *slower* of the two on long
    documents, and it was the one streaming nothing between "start" and a
    result several minutes later.
    """
    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "docx")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        _convert_pdf2docx(decrypted_path, output_file, progress_callback)
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return str(output_file)


def merge_pdfs(input_paths: List[str], output_dir: str, passwords: List[str] = None) -> str:
    """Merge multiple PDFs into a single PDF, in the given order."""
    if not input_paths:
        raise ValueError("No input files provided for merging.")
    if len(input_paths) < 2:
        raise ValueError("Provide at least two PDF files to merge.")

    output_dir_path = Path(output_dir)
    output_file = output_dir_path / f"merged_{uuid.uuid4().hex[:8]}.pdf"

    cleanup_paths: List[str] = []
    try:
        merged = pikepdf.Pdf.new()
        for idx, path in enumerate(input_paths):
            pwd = passwords[idx] if passwords and idx < len(passwords) else None
            decrypted_path, needs_cleanup = _get_decrypted_pdf_path(path, pwd, output_dir_path)
            if needs_cleanup:
                cleanup_paths.append(decrypted_path)
            with pikepdf.open(decrypted_path) as src:
                merged.pages.extend(src.pages)
        merged.save(output_file)
    finally:
        for p in cleanup_paths:
            Path(p).unlink(missing_ok=True)

    return str(output_file)


def add_watermark(
    input_path: str,
    output_dir: str,
    text: str,
    position: str = "diagonal",
    opacity: float = 0.3,
    password: str = None,
) -> str:
    """Stamp a text watermark on every page."""
    import fitz

    if not text or not text.strip():
        raise ValueError("Watermark text cannot be empty.")
    try:
        opacity = float(opacity)
    except (TypeError, ValueError):
        raise ValueError("Opacity must be a number between 0.1 and 1.0.")
    if not 0.05 <= opacity <= 1.0:
        raise ValueError("Opacity must be between 0.1 and 1.0.")
    if position not in ("diagonal", "top", "center", "bottom"):
        raise ValueError("Position must be one of: diagonal, top, center, bottom.")

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        # Use PyMuPDF for the overlay — simpler cross-page-size handling than pikepdf overlays.
        doc = fitz.open(decrypted_path)
        try:
            for page in doc:
                rect = page.rect
                # Pick a font size relative to page width.
                font_size = max(24, int(rect.width / 12))
                color = (0.5, 0.5, 0.5)

                if position == "diagonal":
                    # Diagonal stamp anchored at page center (PyMuPDF rotate= must be a
                    # multiple of 90, so apply the 45° rotation via morph).
                    point = fitz.Point(rect.width / 2, rect.height / 2)
                    page.insert_text(
                        point,
                        text,
                        fontname="helv",
                        fontsize=font_size,
                        color=color,
                        fill_opacity=opacity,
                        morph=(point, fitz.Matrix(1, 1).prerotate(45)),
                    )
                else:
                    if position == "top":
                        y = rect.height * 0.1
                    elif position == "center":
                        y = rect.height / 2
                    else:  # bottom
                        y = rect.height * 0.9
                    # Rough horizontal centering.
                    text_width = font_size * 0.5 * len(text)
                    x = max(10, (rect.width - text_width) / 2)
                    page.insert_text(
                        fitz.Point(x, y),
                        text,
                        fontname="helv",
                        fontsize=font_size,
                        color=color,
                        fill_opacity=opacity,
                    )
            doc.save(str(output_file), garbage=3, deflate=True)
        finally:
            doc.close()
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return str(output_file)


def pdf_to_images_zip(
    input_path: str,
    output_dir: str,
    dpi: int = 150,
    fmt: str = "jpg",
    password: str = None,
) -> dict:
    """Render every PDF page to an image and return a zip."""
    import zipfile
    import fitz

    try:
        dpi = int(dpi)
    except (TypeError, ValueError):
        raise ValueError("DPI must be an integer.")
    if dpi < 50 or dpi > 300:
        raise ValueError("DPI must be between 50 and 300.")

    fmt = (fmt or "jpg").lower()
    if fmt not in ("jpg", "jpeg", "png"):
        raise ValueError("Format must be jpg or png.")
    file_ext = "jpg" if fmt in ("jpg", "jpeg") else "png"

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "zip")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        doc = fitz.open(decrypted_path)
        try:
            page_count = len(doc)
            with zipfile.ZipFile(output_file, "w", compression=zipfile.ZIP_DEFLATED) as zf:
                for i, page in enumerate(doc, start=1):
                    pix = page.get_pixmap(dpi=dpi)
                    if file_ext == "jpg":
                        img_bytes = pix.tobytes("jpeg")
                    else:
                        img_bytes = pix.tobytes("png")
                    arcname = Path(f"{original_stem(input_file)}_page_{i:03d}.{file_ext}").name
                    zf.writestr(arcname, img_bytes)
                    pix = None
        finally:
            doc.close()
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return {"output_path": str(output_file), "page_count": page_count}


def sign_pdf(
    input_path: str,
    signature_image_path: str,
    output_dir: str,
    page: int = 1,
    x: float = 0.65,
    y: float = 0.85,
    width: float = 0.2,
    password: str = None,
) -> str:
    """Stamp a signature image onto the chosen page.

    x, y, width are normalized (0-1) relative to page size. (x, y) is the top-left
    of the signature box.
    """
    import fitz

    try:
        page = int(page)
        x = float(x)
        y = float(y)
        width = float(width)
    except (TypeError, ValueError):
        raise ValueError("page must be an integer; x, y, width must be numbers.")
    if page < 1:
        raise ValueError("Page number must be >= 1.")
    if not (0.0 <= x <= 1.0 and 0.0 <= y <= 1.0):
        raise ValueError("x and y must be between 0 and 1.")
    if not (0.05 <= width <= 1.0):
        raise ValueError("width must be between 0.05 and 1.0.")
    if not Path(signature_image_path).exists():
        raise ValueError("Signature image not found.")

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        doc = fitz.open(decrypted_path)
        try:
            if page > len(doc):
                raise ValueError(f"Page {page} exceeds document page count ({len(doc)}).")

            target_page = doc[page - 1]
            rect = target_page.rect
            box_w = rect.width * width
            # Estimate height from image aspect ratio so the signature isn't squashed.
            try:
                from PIL import Image as PILImage
                with PILImage.open(signature_image_path) as im:
                    img_w, img_h = im.size
                aspect = img_h / img_w if img_w else 0.4
            except Exception:
                aspect = 0.4
            box_h = box_w * aspect

            x0 = rect.width * x
            y0 = rect.height * y
            x1 = min(rect.width, x0 + box_w)
            y1 = min(rect.height, y0 + box_h)
            target_rect = fitz.Rect(x0, y0, x1, y1)

            target_page.insert_image(target_rect, filename=signature_image_path, keep_proportion=True)
            doc.save(str(output_file), garbage=3, deflate=True)
        finally:
            doc.close()
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return str(output_file)

def merge_docx_files(input_files: list, output_file: str) -> None:
    """Merges multiple DOCX files into one, inserting page breaks between them."""
    if not input_files:
        raise ValueError("No input files provided for merging.")

    if Document_docx is None or Composer is None:
        raise ImportError("docxcompose and python-docx are required for DOCX merging")

    master = Document_docx(input_files[0])
    composer = Composer(master)

    for docx_path in input_files[1:]:
        master.add_page_break()
        composer.append(Document_docx(str(docx_path)))

    composer.save(output_file)

def rotate_pdf(input_path: str, output_dir: str, angle: int, pages: str = None, password: str = None) -> str:
    """Rotate PDF pages by specified angle (90, 180, 270).

    Args:
        input_path: Path to input PDF
        output_dir: Directory to save output PDF
        angle: Rotation angle (90, 180, 270, or -90, -180, -270)
        pages: Page selection string (e.g., '1,3-5', 'all'). If None, rotates all pages.
        password: PDF password if encrypted

    Returns:
        Path to rotated PDF file
    """
    if angle not in (90, 180, 270, -90, -180, -270):
        raise ValueError("Angle must be 90, 180, 270, -90, -180, or -270 degrees.")

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        with pikepdf.open(decrypted_path) as pdf:
            total_pages = len(pdf.pages)

            # Parse page selection (default to all if not specified)
            if pages is None:
                selected_indices = list(range(total_pages))
            else:
                selected_indices = _parse_page_selection(pages, total_pages)

            # Rotate selected pages
            for idx in selected_indices:
                page = pdf.pages[idx]
                # pikepdf uses /Rotate key: 0, 90, 180, 270 (only these values)
                current_rotation = int(page.get('/Rotate', 0))
                new_rotation = (current_rotation + angle) % 360
                page['/Rotate'] = new_rotation

            pdf.save(output_file)
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return str(output_file)


def _safe_progress(progress_callback, done: int, total: int) -> None:
    if progress_callback:
        try:
            progress_callback(done, total)
        except Exception:
            logger.debug("Progress callback failed", exc_info=True)


def _pdf_to_word_paddle_impl(decrypted_path: str, temp_dir: Path, output_file: Path,
                             progress_callback=None) -> None:
    """PPStructure layout-recovery conversion (Paddle backend, x86 only)."""
    import fitz

    # Deferred imports for utility functions that depend on paddleocr
    try:
        from paddleocr import save_structure_res
        from paddleocr.ppstructure.recovery.recovery_to_doc import sorted_layout_boxes, convert_info_docx
    except ImportError:
        raise ImportError("PaddleOCR sub-modules could not be loaded. Please check your installation.")

    # Use cached PaddleOCR engine instead of re-initializing
    table_engine = get_paddle_engine()

    doc = fitz.open(decrypted_path)
    total_pages = len(doc)
    logger.info("Opened PDF with %d pages", total_pages)
    docx_files = []

    _safe_progress(progress_callback, 0, total_pages)

    for i, page in enumerate(doc):
        # 200 DPI is a good balance between speed and quality for OCR
        img = _render_page_bgr(page)

        # Run inference
        result = table_engine(img)

        # Save structure result (images, excels)
        page_name = f"page_{i}"
        save_structure_res(result, str(temp_dir), page_name)

        # Convert to DOCX using recovery module
        h, w, _ = img.shape
        res = sorted_layout_boxes(result, w)
        convert_info_docx(img, res, str(temp_dir), page_name)

        # The docx is saved as {page_name}_ocr.docx in temp_dir
        expected_docx = temp_dir / f"{page_name}_ocr.docx"

        if expected_docx.exists():
            docx_files.append(expected_docx)
        else:
            logger.warning("Could not find recovered docx for page %d", i)

        _safe_progress(progress_callback, i + 1, total_pages)

    if not docx_files:
        raise Exception("No pages were successfully converted using AI engine.")

    # Merge recovered per-page DOCX files with page breaks between them
    merge_docx_files([str(f) for f in docx_files], str(output_file))
    doc.close()


def _split_into_columns(entries: List[tuple], page_width: Optional[float],
                        min_gap_frac: float = 0.06) -> List[List[tuple]]:
    """Split OCR text fragments into left-to-right column bands using a
    whitespace-gap projection onto the x-axis.

    Prevents fragments from different columns (common on multi-column
    resumes) from being merged into the same reading-order line just
    because their y-centers happen to be close.
    """
    if not entries or not page_width:
        return [entries]

    spans = sorted((e[3], e[4]) for e in entries)  # (x0, x1)
    merged: List[List[float]] = []
    for x0, x1 in spans:
        if merged and x0 <= merged[-1][1]:
            merged[-1][1] = max(merged[-1][1], x1)
        else:
            merged.append([x0, x1])

    min_gap = page_width * min_gap_frac
    boundaries = [a[1] for a, b in zip(merged, merged[1:]) if b[0] - a[1] >= min_gap]
    if not boundaries:
        return [entries]

    columns: List[List[tuple]] = [[] for _ in range(len(boundaries) + 1)]
    for entry in entries:
        x_center = (entry[3] + entry[4]) / 2
        idx = sum(1 for b in boundaries if x_center > b)
        columns[idx].append(entry)
    return [c for c in columns if c]


def _group_ocr_lines(items: List[dict], page_width: Optional[float] = None) -> List[str]:
    """Group OCR text snippets into reading-order lines using their bboxes.

    Text is first split into left-to-right columns (`_split_into_columns`)
    so multi-column layouts don't interleave. Within each column, snippets
    whose vertical centers fall within roughly half a text height of each
    other are treated as one line and joined left-to-right - but only if
    they don't horizontally overlap with what's already in that line,
    since overlapping x-ranges can't belong to the same visual row (this
    keeps a new bullet's left-margin start from being absorbed into a
    still-open previous row just because their y-centers are close).

    Known limitation: if a single detection box spans multiple
    tightly-kerned words (e.g. "SeniorAnalyticsConsultant"), the space
    between them is lost here too - RapidOCR's recognizer has no explicit
    space class, and it doesn't expose sub-box geometry this function could
    use to split it back apart. A dictionary-segmentation package
    (`wordninja`) was evaluated for this, but its legacy sdist-only build
    fails under current setuptools (>= ~68, which is what File-Forge's
    Docker build installs),
    so this is intentionally left unfixed rather than shipping a fragile
    dependency or a heuristic that can silently produce wrong text. This
    only affects pages that still require true OCR (no embedded text layer
    at all) - see `pdf_to_word_ai`'s text-layer routing, which is what
    prevents OCR from running on normal, text-based PDFs in the first place.
    """
    entries = []
    for item in items:
        text = (item.get("text") or "").strip()
        if not text:
            continue
        bbox = item.get("bbox") or []
        try:
            xs = [float(pt[0]) for pt in bbox]
            ys = [float(pt[1]) for pt in bbox]
            x0, x1, y0, y1 = min(xs), max(xs), min(ys), max(ys)
        except (TypeError, ValueError, IndexError):
            x0, x1, y0, y1 = 0.0, 0.0, float(len(entries)), float(len(entries))
        entries.append((y0, (y0 + y1) / 2, y1 - y0, x0, x1, text))

    lines_out: List[str] = []
    for column_entries in _split_into_columns(entries, page_width):
        column_entries = sorted(column_entries, key=lambda e: (e[0], e[3]))

        lines: List[List[tuple]] = []
        for entry in column_entries:
            if lines:
                prev = lines[-1]
                prev_center = sum(e[1] for e in prev) / len(prev)
                prev_max_x1 = max(e[4] for e in prev)
                tolerance = max(entry[2], 1.0) * 0.6
                overlap_tolerance = max(entry[2], 1.0) * 0.15
                same_row = abs(entry[1] - prev_center) <= tolerance
                overlaps_existing = entry[3] < prev_max_x1 - overlap_tolerance
                if same_row and not overlaps_existing:
                    prev.append(entry)
                    continue
            lines.append([entry])

        for line in lines:
            lines_out.append(
                " ".join(e[-1] for e in sorted(line, key=lambda e: e[3]))
            )

    return lines_out


def _pdf_to_word_ocr_fallback(engine, decrypted_path: str, output_file: Path,
                              progress_callback=None) -> None:
    """Plain-OCR conversion: recognize text per page and write it to a DOCX.

    Used when the configured backend (e.g. RapidOCR on ARM64) has no layout
    recovery. Preserves reading order but not tables/columns/figures.
    """
    import fitz

    if Document_docx is None:
        raise ImportError("python-docx is required for OCR-based PDF to Word conversion")

    doc = fitz.open(decrypted_path)
    total_pages = len(doc)
    logger.info("Opened PDF with %d pages (OCR fallback via %s)", total_pages, engine.name)

    word_doc = Document_docx()
    any_text = False

    _safe_progress(progress_callback, 0, total_pages)

    for i, page in enumerate(doc):
        img = _render_page_bgr(page)
        items = engine.recognize(img)

        if i > 0:
            word_doc.add_page_break()
        for line in _group_ocr_lines(items, page_width=page.rect.width):
            word_doc.add_paragraph(line)
            any_text = True

        _safe_progress(progress_callback, i + 1, total_pages)

    doc.close()

    if not any_text:
        raise ValueError("No text could be recognized in this PDF.")

    word_doc.save(str(output_file))


def _pdf_to_word_hybrid_impl(engine, decrypted_path: str, output_file: Path,
                             report: dict, progress_callback=None) -> None:
    """Mixed-content PDF: extract native text where a page already has a
    usable text layer, and only rasterize+OCR the pages that don't."""
    import fitz

    if Document_docx is None:
        raise ImportError("python-docx is required for OCR-based PDF to Word conversion")

    doc = fitz.open(decrypted_path)
    total_pages = len(doc)
    logger.info(
        "Opened PDF with %d pages (hybrid: native text + OCR via %s for scanned pages)",
        total_pages, engine.name,
    )

    word_doc = Document_docx()
    any_text = False

    _safe_progress(progress_callback, 0, total_pages)

    for i, page in enumerate(doc):
        if i > 0:
            word_doc.add_page_break()

        if report["pages_with_text"][i]:
            for para in page.get_text("text").split("\n"):
                para = para.strip()
                if para:
                    word_doc.add_paragraph(para)
                    any_text = True
        else:
            img = _render_page_bgr(page)
            items = engine.recognize(img)
            for line in _group_ocr_lines(items, page_width=page.rect.width):
                word_doc.add_paragraph(line)
                any_text = True

        _safe_progress(progress_callback, i + 1, total_pages)

    doc.close()

    if not any_text:
        raise ValueError("No text could be recognized in this PDF.")

    word_doc.save(str(output_file))


def pdf_to_word_ai(input_path: str, output_dir: str, password: str = None,
                   progress_callback=None, method_callback=None) -> str:
    """Converts PDF to DOCX, preferring the PDF's own text layer over OCR.

    Rasterize+OCR is expensive and lossy (dropped spacing, character
    misreads, no font/layout preservation), so it should only ever run on
    pages that don't already have usable extractable text - not as the
    default outcome of checking "AI Layout Recovery" on a normal PDF.

    Routing (in order):
        1. The document already has a usable embedded text layer (>=80% of
           pages qualify) -> standard pdf2docx conversion. OCR is never
           invoked, regardless of the caller's AI intent.
        2. No usable text layer, configured engine supports layout recovery
           (paddle/PPStructure, x86 only) -> layout-aware OCR conversion.
        3. No usable text layer, but some pages have text and some don't,
           and the engine lacks layout support -> hybrid conversion: native
           text per page where present, OCR only the pages that need it.
        4. No text anywhere, no layout support -> full-page OCR fallback.

    Args:
        progress_callback: Optional callable(page_done, total_pages) invoked
            after each page is processed, for streaming progress to clients.
        method_callback: Optional callable(str) invoked once with one of
            "text_layer" | "paddle_layout" | "ocr_hybrid" | "ocr_fallback"
            as soon as the routing decision is made, so callers can report
            an accurate result message.
    """
    import fitz

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "docx")

    # Create a temp directory for intermediate files (incl. decrypted copies)
    temp_dir = Path(output_dir) / f"temp_{original_stem(input_file)}"
    temp_dir.mkdir(exist_ok=True)

    # Handle encrypted PDFs
    logger.debug("Checking encryption...")
    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password, temp_dir)
    logger.debug("Using path: %s, needs_cleanup: %s", decrypted_path, needs_cleanup)

    try:
        probe = fitz.open(decrypted_path)
        try:
            report = _inspect_text_layer(probe)
        finally:
            probe.close()

        if report["has_usable_text_layer"]:
            logger.info(
                "pdf_to_word_ai: %s already has a usable embedded text layer "
                "(%d/%d pages) - using standard text-based conversion instead "
                "of rasterize+OCR.",
                input_file.name, sum(report["pages_with_text"]), report["total_pages"],
            )
            if method_callback:
                method_callback("text_layer")
            _safe_progress(progress_callback, 0, report["total_pages"])
            _convert_pdf2docx(decrypted_path, output_file)
            _safe_progress(progress_callback, report["total_pages"], report["total_pages"])
            return str(output_file)

        # Only reachable when pages genuinely lack usable text (scanned /
        # image-only) - only now is it worth spinning up the OCR engine.
        from scripts.ocr_engine import get_ocr_engine

        engine = get_ocr_engine()
        if engine is None:
            raise ValueError(
                "This PDF appears to be scanned/image-based, and AI Layout "
                "Recovery is disabled on this server, so it can't be "
                "converted. (A PDF with selectable text would convert fine "
                "without AI.)"
            )

        logger.info("Starting AI conversion for: %s (backend=%s)", input_path, engine.name)

        if engine.supports_layout:
            if method_callback:
                method_callback("paddle_layout")
            _pdf_to_word_paddle_impl(decrypted_path, temp_dir, output_file, progress_callback)
        elif report["fraction_with_text"] > 0.0:
            if method_callback:
                method_callback("ocr_hybrid")
            _pdf_to_word_hybrid_impl(engine, decrypted_path, output_file, report, progress_callback)
        else:
            if method_callback:
                method_callback("ocr_fallback")
            _pdf_to_word_ocr_fallback(engine, decrypted_path, output_file, progress_callback)
    finally:
        # Cleanup
        if temp_dir.exists():
            try:
                shutil.rmtree(temp_dir)
            except Exception:
                logger.warning("Could not fully clean up %s", temp_dir)

    return str(output_file)


# Backwards-compatible alias (pre-ARM-migration name)
pdf_to_word_paddle = pdf_to_word_ai


# ──────────────────────────────────────────────
# Feature #53: Protect PDF (Add Password)
# ──────────────────────────────────────────────

def protect_pdf(
    input_path: str,
    output_dir: str,
    user_password: str,
    owner_password: str = None,
    allow_print: bool = True,
    allow_copy: bool = False,
    allow_edit: bool = False,
    password: str = None,
) -> str:
    """Add password protection and access restrictions to a PDF.

    Args:
        input_path: Path to input PDF.
        output_dir: Directory to save protected PDF.
        user_password: Password required to open the document.
        owner_password: Password granting unrestricted access (defaults to user_password).
        allow_print: Whether to allow printing.
        allow_copy: Whether to allow text copying.
        allow_edit: Whether to allow editing.
        password: Existing password if the PDF is already encrypted.

    Returns:
        Path to the protected PDF.
    """
    if not user_password:
        raise ValueError("User password cannot be empty.")

    owner_pw = owner_password or user_password

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        permissions = pikepdf.Permissions(
            print_lowres=allow_print,
            print_highres=allow_print,
            extract=allow_copy,
            modify_annotation=allow_edit,
            modify_assembly=allow_edit,
            modify_form=allow_edit,
            modify_other=allow_edit,
            accessibility=True,  # Always allow accessibility
        )
        encryption = pikepdf.Encryption(
            user=user_password,
            owner=owner_pw,
            allow=permissions,
        )
        with pikepdf.open(decrypted_path) as pdf:
            pdf.save(output_file, encryption=encryption)
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return str(output_file)


# ──────────────────────────────────────────────
# Feature #54: Image to PDF
# ──────────────────────────────────────────────

# EXIF Orientation tag id (same value used in image_utils).
_EXIF_ORIENTATION_TAG = 0x0112


def _oriented_for_pdf(img_path: str, tmp_dir: str, tmp_files: list) -> tuple:
    """Return a path whose pixels are correctly oriented, plus its (width, height).

    reportlab's ``drawImage`` reads raw pixels and ignores the EXIF Orientation
    tag, so a portrait phone photo (stored landscape + ``Orientation=6``) would be
    embedded **sideways** in the generated PDF. That is a concrete fidelity
    violation for images→PDF.

    Fast path (the common case): images with no rotation — Orientation absent, ``0``
    or ``1`` — are returned **unchanged**, so already-upright inputs keep the exact,
    un-re-encoded embedding the original code produced (byte-identical output).

    Only images that actually carry a non-identity Orientation get their rotation
    baked into a temporary copy via ``ImageOps.exif_transpose`` (a lossless pixel
    transpose). JPEG sources are re-saved as high-quality JPEG to keep the PDF
    small; everything else is saved as lossless PNG (which also preserves alpha).
    Any temp file created is appended to ``tmp_files`` for the caller to clean up.
    """
    from PIL import Image as _PILImage, ImageOps as _ImageOps

    with _PILImage.open(img_path) as img:
        try:
            orientation = img.getexif().get(_EXIF_ORIENTATION_TAG, 1)
        except Exception:
            orientation = 1

        if orientation in (0, 1):
            return img_path, img.size[0], img.size[1]

        transposed = _ImageOps.exif_transpose(img)
        suffix = Path(img_path).suffix.lower()
        if suffix in (".jpg", ".jpeg"):
            tmp = Path(tmp_dir) / f"_oriented_{uuid.uuid4().hex[:8]}.jpg"
            save_img = transposed if transposed.mode in ("RGB", "L") else transposed.convert("RGB")
            save_img.save(tmp, "JPEG", quality=95, optimize=True)
        else:
            tmp = Path(tmp_dir) / f"_oriented_{uuid.uuid4().hex[:8]}.png"
            transposed.save(tmp, "PNG")
        tmp_files.append(tmp)
        w, h = transposed.size
        return str(tmp), w, h


def images_to_pdf(
    input_paths: List[str],
    output_dir: str,
    page_size: str = "A4",
    fit_mode: str = "fit",
    margin_pt: int = 36,
) -> str:
    """Convert one or more images into a single multi-page PDF.

    Args:
        input_paths: Ordered list of image file paths.
        output_dir: Directory to save the PDF.
        page_size: 'A4', 'Letter', or 'auto' (use image dimensions).
        fit_mode: 'fit' (scale to page), 'stretch', or 'original'.
        margin_pt: Page margin in points (default 36 = 0.5 inch).

    Returns:
        Path to the created PDF.
    """
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.pagesizes import A4, LETTER
    from PIL import Image as PILImage

    if not input_paths:
        raise ValueError("At least one image file is required.")

    PAGE_SIZES = {"a4": A4, "letter": LETTER}
    page_size_key = page_size.lower()

    output_name = f"images_to_pdf_{uuid.uuid4().hex[:8]}.pdf"
    output_file = Path(output_dir) / output_name

    # Use a temp canvas path, we'll rebuild after getting page sizes
    c = rl_canvas.Canvas(str(output_file))

    # Orientation-normalized temp copies (only created for rotated inputs); removed
    # after the canvas is saved since reportlab reads each image at drawImage time.
    tmp_files: list = []
    try:
        for img_path in input_paths:
            # reportlab ignores EXIF orientation, so bake it in for rotated inputs.
            draw_path, img_w, img_h = _oriented_for_pdf(img_path, output_dir, tmp_files)
            if img_w == 0 or img_h == 0:
                raise ValueError(f"Image {img_path} has zero-dimension (width={img_w}, height={img_h}).")

            if page_size_key == "auto":
                pw, ph = float(img_w), float(img_h)
            else:
                pw, ph = PAGE_SIZES.get(page_size_key, A4)

            c.setPageSize((pw, ph))

            available_w = pw - 2 * margin_pt
            available_h = ph - 2 * margin_pt

            if fit_mode == "original":
                draw_w = min(float(img_w), available_w)
                draw_h = draw_w * (img_h / img_w)
            else:
                scale = min(available_w / img_w, available_h / img_h)
                draw_w = img_w * scale
                draw_h = img_h * scale

            x = margin_pt + (available_w - draw_w) / 2
            y = margin_pt + (available_h - draw_h) / 2

            c.drawImage(draw_path, x, y, width=draw_w, height=draw_h, preserveAspectRatio=True)
            c.showPage()

        c.save()
    finally:
        for tmp in tmp_files:
            try:
                os.remove(tmp)
            except OSError:
                pass
    return str(output_file)


# ──────────────────────────────────────────────
# Feature #55: Word to PDF
# ──────────────────────────────────────────────

def word_to_pdf(input_path: str, output_dir: str) -> str:
    """Convert a DOCX/DOC file to PDF.

    Tries LibreOffice first; falls back to a pure-Python
    python-docx + reportlab converter.

    Args:
        input_path: Path to the DOCX/DOC file.
        output_dir: Directory to save the PDF.

    Returns:
        Path to the converted PDF.
    """
    import subprocess

    input_file = Path(input_path)
    suffix = input_file.suffix.lower()
    if suffix not in (".docx", ".doc", ".odt", ".rtf"):
        raise ValueError("Input must be a .docx, .doc, .odt, or .rtf file.")

    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    # ── Try LibreOffice first ──────────────────────────────────
    try:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf",
             "--outdir", str(output_dir), str(input_file)],
            capture_output=True, text=True, timeout=120,
        )
        # LibreOffice writes "<input stem>.pdf" itself; rename to our branded
        # name rather than checking for a path it never produces.
        libreoffice_output = Path(output_dir) / f"{input_file.stem}.pdf"
        if result.returncode == 0 and libreoffice_output.exists():
            if libreoffice_output != output_file:
                libreoffice_output.replace(output_file)
            return str(output_file)
        # Reached LibreOffice but it refused the document. Logged, because the
        # fallback below silently produces a lower-fidelity result and the
        # reason is otherwise lost.
        logger.warning(
            "LibreOffice could not convert %s (exit %s): %s",
            input_file.name, result.returncode, (result.stderr or "").strip()[:300],
        )
    except FileNotFoundError:
        # LibreOffice isn't installed at all — the state the Oracle VM was found
        # in on 2026-07-21. Every .doc/.odt/.rtf conversion fails outright and
        # .docx quietly drops to the low-fidelity path; worth saying loudly
        # once per conversion rather than swallowing.
        logger.error(
            "LibreOffice is not installed on this host, so %s conversions fall back "
            "to the low-fidelity pure-Python path (and non-.docx input fails "
            "outright). See docs/fix-libreoffice-libgl-oracle.md.", suffix,
        )
    except subprocess.TimeoutExpired:
        logger.warning("LibreOffice timed out converting %s", input_file.name)
    except Exception:
        logger.warning("LibreOffice conversion of %s failed", input_file.name, exc_info=True)

    # ── Pure-Python fallback (DOCX only) ──────────────────────
    if suffix != ".docx":
        raise RuntimeError(
            "LibreOffice conversion failed. Pure-Python fallback only supports .docx files."
        )

    try:
        from docx import Document as DocxDocument
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.lib.units import inch
    except ImportError as e:
        raise RuntimeError(f"python-docx or reportlab is required for DOCX→PDF fallback: {e}")

    docx = DocxDocument(str(input_file))
    doc = SimpleDocTemplate(str(output_file), pagesize=A4,
                            leftMargin=inch, rightMargin=inch,
                            topMargin=inch, bottomMargin=inch)

    styles = getSampleStyleSheet()
    heading_style = ParagraphStyle("Heading", parent=styles["Heading1"], spaceAfter=6)
    body_style = ParagraphStyle("Body", parent=styles["Normal"], fontSize=11,
                                leading=15, spaceAfter=4)

    def _escape(text: str) -> str:
        return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    elements = []
    for para in docx.paragraphs:
        text = para.text.strip()
        if not text:
            elements.append(Spacer(1, 6))
            continue
        style = heading_style if para.style.name.startswith("Heading") else body_style
        elements.append(Paragraph(_escape(text), style))

    for table in docx.tables:
        data = []
        for row in table.rows:
            data.append([_escape(cell.text) for cell in row.cells])
        if data:
            tbl = Table(data, repeatRows=1)
            tbl.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
            ]))
            elements.append(tbl)
            elements.append(Spacer(1, 8))

    if not elements:
        elements.append(Paragraph("(Empty document)", body_style))

    doc.build(elements)
    return str(output_file)


def word_to_pptx(input_path: str, output_dir: str, dpi: int = 150) -> str:
    """Convert a DOCX/DOC file to a PowerPoint presentation.

    Renders the document to PDF first (via word_to_pdf), then rasterizes each
    resulting page as a slide image (via pdf_to_pptx) — the same non-editable,
    image-slide approach pdf_to_pptx already uses for PDF input, so a page's
    layout survives the trip without needing OCR.

    Args:
        input_path: Path to the DOCX/DOC/ODT/RTF file.
        output_dir: Directory to save the PPTX file.
        dpi: Rendering resolution for each slide image.

    Returns:
        Path to the converted PPTX file.
    """
    input_file = Path(input_path)
    output_dir_path = Path(output_dir)
    # word_to_pdf names its output "<original>_forgefiles.org.pdf"; renaming it
    # to a UUID-prefixed temp name before handing it to pdf_to_pptx lets that
    # function's own branding recover the true original stem, instead of
    # branding the already-branded name a second time. input_file may itself
    # still carry an upload UUID prefix, so strip that first or it survives
    # as part of the "stem" pdf_to_pptx thinks is the original name.
    temp_pdf_path = output_dir_path / f"{uuid.uuid4()}_{original_stem(input_file)}.pdf"
    try:
        rendered_pdf = word_to_pdf(str(input_file), str(output_dir_path))
        Path(rendered_pdf).replace(temp_pdf_path)
        return pdf_to_pptx(str(temp_pdf_path), str(output_dir_path), dpi=dpi)
    finally:
        if temp_pdf_path.exists():
            temp_pdf_path.unlink(missing_ok=True)


# ──────────────────────────────────────────────
# Feature #56: PDF to Excel (Table Extraction)
# ──────────────────────────────────────────────

def _extract_borderless_tables(page) -> list:
    """Recover whitespace-aligned (borderless) tables the default 'lines'
    strategy cannot see.

    ``page.find_tables()`` defaults to the "lines" strategy, which relies on
    drawn ruling lines/rectangles. Tables that separate columns by whitespace
    alone (common in financial statements, invoices, bank statements) carry no
    such geometry and are therefore missed entirely — the page contributes no
    tables and its data is lost. PyMuPDF's own text-position strategy
    (``strategy="text"``) reconstructs virtual column/row boundaries from word
    alignment and is the documented remedy for borderless tables.

    Called only for pages where the "lines" pass found nothing, so it never
    alters an already-detected table. Returns a list of cleaned row-lists
    (fully-empty spacer rows trimmed). To avoid mangling ordinary prose into a
    spurious table, only genuine grids (>= 2 columns and >= 2 non-empty rows)
    are returned; anything else yields ``[]`` so the page still falls through to
    the raw-text fallback sheet exactly as before.
    """
    try:
        found = page.find_tables(
            vertical_strategy="text",
            horizontal_strategy="text",
        )
    except Exception:
        # Text-strategy detection is best-effort; never let it break extraction.
        return []

    accepted = []
    for table in getattr(found, "tables", []) or []:
        if getattr(table, "col_count", 0) < 2:
            continue
        rows = [
            [("" if cell is None else cell) for cell in row]
            for row in table.extract()
            if any(cell is not None and str(cell).strip() != "" for cell in row)
        ]
        if len(rows) < 2:
            continue
        accepted.append(rows)
    return accepted


def pdf_to_excel(input_path: str, output_dir: str, password: str = None) -> dict:
    """Extract tables from a PDF and save to an Excel workbook.

    Uses PyMuPDF's find_tables (default "lines" strategy) for extraction, with a
    text-position fallback for borderless tables on pages the "lines" strategy
    can't see (see :func:`_extract_borderless_tables`).

    Args:
        input_path: Path to the input PDF.
        output_dir: Directory to save the Excel file.
        password: PDF password if encrypted.

    Returns:
        dict with output_path and tables_found count.
    """
    import openpyxl

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "xlsx")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        doc = fitz.open(decrypted_path)
        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # Remove default empty sheet
        tables_found = 0

        for page_num, page in enumerate(doc, start=1):
            tables = page.find_tables()
            if tables and tables.tables:
                for tbl_idx, table in enumerate(tables.tables):
                    tables_found += 1
                    sheet_name = f"P{page_num}_T{tbl_idx + 1}"[:31]  # Excel max 31 chars
                    ws = wb.create_sheet(title=sheet_name)
                    for row in table.extract():
                        ws.append([cell if cell is not None else "" for cell in row])
            else:
                # No ruled table on this page — try to recover a borderless one
                # before giving up (additive; leaves ruled-table output intact).
                for tbl_idx, rows in enumerate(_extract_borderless_tables(page)):
                    tables_found += 1
                    sheet_name = f"P{page_num}_T{tbl_idx + 1}"[:31]  # Excel max 31 chars
                    ws = wb.create_sheet(title=sheet_name)
                    for row in rows:
                        ws.append([cell if cell is not None else "" for cell in row])

        doc.close()

        if tables_found == 0:
            # No tables found — extract raw text into a single sheet as fallback
            doc2 = fitz.open(decrypted_path)
            ws = wb.create_sheet(title="Text Content")
            ws.append(["Page", "Text"])
            for page_num, page in enumerate(doc2, start=1):
                text = page.get_text().strip()
                if text:
                    ws.append([page_num, text])
            doc2.close()

        wb.save(output_file)
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return {"output_path": str(output_file), "tables_found": tables_found}


# ──────────────────────────────────────────────
# Feature #57: PDF to PowerPoint
# ──────────────────────────────────────────────

def pdf_to_pptx(
    input_path: str,
    output_dir: str,
    dpi: int = 150,
    password: str = None,
) -> str:
    """Convert each PDF page to a slide in a PowerPoint presentation.

    Each page is rendered as a high-res image and used as the slide background.

    Args:
        input_path: Path to input PDF.
        output_dir: Directory to save the PPTX file.
        dpi: Rendering resolution.
        password: PDF password if encrypted.

    Returns:
        Path to the created PPTX file.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    import io

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pptx")

    try:
        dpi = int(dpi)
        if dpi < 72 or dpi > 300:
            raise ValueError("DPI must be between 72 and 300.")
    except (TypeError, ValueError):
        raise ValueError("DPI must be an integer between 72 and 300.")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        doc = fitz.open(decrypted_path)
        prs = Presentation()

        # PowerPoint supports only ONE slide size for the whole deck (a hard limit
        # of the .pptx format, not just python-pptx). The previous code set
        # prs.slide_width/height *inside* the per-page loop, so on a mixed-page-size
        # PDF only the LAST page's size survived: every earlier slide kept a
        # full-page image sized to its own page but sat on a differently-sized
        # canvas, clipping or off-setting the page. Instead, size the deck once to
        # the bounding box of all pages (max width, max height in EMU) so no page
        # can overflow, then place each page's image at its native size, centered —
        # preserving each page's aspect ratio (no stretching) and clipping nothing.
        # For a uniform-size PDF the bounding box equals the page size, so every
        # image fills its slide exactly at (0, 0) — identical to the old output.
        EMU_PER_PT = 914400 / 72.0
        page_count = len(doc)
        page_sizes_emu = []
        for i in range(page_count):
            rect = doc[i].rect
            page_sizes_emu.append(
                (int(rect.width * EMU_PER_PT), int(rect.height * EMU_PER_PT))
            )

        prs.slide_width = max(
            (w for w, _ in page_sizes_emu), default=int(8.5 * 914400)
        )
        prs.slide_height = max(
            (h for _, h in page_sizes_emu), default=int(11 * 914400)
        )

        for i in range(page_count):
            page = doc[i]
            page_w, page_h = page_sizes_emu[i]

            # Render page to image bytes (unchanged — lossless PNG at the DPI)
            pix = page.get_pixmap(dpi=dpi)
            img_bytes = io.BytesIO(pix.tobytes("png"))

            blank_layout = prs.slide_layouts[6]  # blank layout
            slide = prs.slides.add_slide(blank_layout)

            # Center the page image on the shared canvas at its native size —
            # aspect preserved, nothing clipped. left/top are 0 for uniform decks.
            left = (prs.slide_width - page_w) // 2
            top = (prs.slide_height - page_h) // 2
            slide.shapes.add_picture(
                img_bytes, left, top,
                width=page_w,
                height=page_h,
            )
            pix = None

        doc.close()
        prs.save(str(output_file))
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return str(output_file)


# ──────────────────────────────────────────────
# Feature #58: Extract Text from PDF
# ──────────────────────────────────────────────

def extract_text_from_pdf(
    input_path: str,
    output_dir: str,
    preserve_layout: bool = False,
    password: str = None,
) -> dict:
    """Extract all text content from a PDF to a .txt file.

    Any page whose native text layer is genuinely empty (scanned/image-only)
    is rasterized and OCR'd via the configured backend, so this - unlike a
    plain text-layer read - actually delivers on "batch OCR" for the premium
    tier (see batch_ocr.py) instead of silently skipping scanned pages. Pages
    that already have native text (however little) keep it as-is; OCR is
    never used to "improve" text that's already there, and the OCR engine is
    only loaded when at least one page is genuinely empty.

    Args:
        input_path: Path to input PDF.
        output_dir: Directory to save the text file.
        preserve_layout: If True, use 'blocks' layout for natively-extracted
            pages; otherwise plain text. OCR'd pages always use the OCR
            engine's recognized line order regardless of this flag.
        password: PDF password if encrypted.

    Returns:
        dict with output_path and page_count.
    """
    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "txt")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        doc = fitz.open(decrypted_path)
        page_count = len(doc)

        native_page_text = []
        for page in doc:
            if preserve_layout:
                blocks = page.get_text("blocks")
                text = "\n".join(b[4].strip() for b in blocks if b[4].strip())
            else:
                text = page.get_text().strip()
            native_page_text.append(text)

        engine = None
        if any(not text for text in native_page_text):
            from scripts.ocr_engine import get_ocr_engine
            engine = get_ocr_engine()

        all_text = []
        for i, page in enumerate(doc):
            page_text = native_page_text[i]

            if not page_text and engine is not None:
                img = _render_page_bgr(page)
                items = engine.recognize(img)
                page_text = "\n".join(item["text"] for item in items if item.get("text")).strip()

            if page_text:
                all_text.append(f"--- Page {i + 1} ---\n{page_text}")

        doc.close()

        full_text = "\n\n".join(all_text) if all_text else "(No text found in document)"
        output_file.write_text(full_text, encoding="utf-8")
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return {"output_path": str(output_file), "page_count": page_count}


# ──────────────────────────────────────────────
# Feature #59: Organize PDF (Reorder/Delete/Duplicate)
# ──────────────────────────────────────────────

def organize_pdf(
    input_path: str,
    output_dir: str,
    page_order: List[int],
    password: str = None,
) -> str:
    """Reorder, delete, or duplicate PDF pages.

    Args:
        input_path: Path to input PDF.
        output_dir: Directory to save the reorganized PDF.
        page_order: 1-based list of page numbers in desired order.
                    Repeat a number to duplicate; omit to delete.
                    Example: [3, 1, 2, 1] — puts page 3 first, duplicates page 1.
        password: PDF password if encrypted.

    Returns:
        Path to the reorganized PDF.
    """
    if not page_order:
        raise ValueError("page_order cannot be empty.")

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        with pikepdf.open(decrypted_path) as pdf:
            total = len(pdf.pages)
            # Validate all page numbers
            for pnum in page_order:
                if not isinstance(pnum, int) or pnum < 1 or pnum > total:
                    raise ValueError(
                        f"Page number {pnum} is out of range (document has {total} pages)."
                    )

            new_pdf = pikepdf.Pdf.new()
            for pnum in page_order:
                new_pdf.pages.append(pdf.pages[pnum - 1])
            new_pdf.save(output_file)
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return str(output_file)


# ──────────────────────────────────────────────
# Feature #60: Add Page Numbers to PDF
# ──────────────────────────────────────────────

def add_page_numbers(
    input_path: str,
    output_dir: str,
    position: str = "bottom-center",
    start_number: int = 1,
    font_size: int = 12,
    skip_first: int = 0,
    fmt: str = "decimal",
    password: str = None,
) -> str:
    """Insert page numbers on each page of a PDF.

    Args:
        input_path: Path to input PDF.
        output_dir: Directory to save the numbered PDF.
        position: One of 'bottom-center', 'bottom-left', 'bottom-right',
                  'top-center', 'top-left', 'top-right'.
        start_number: First page number to use.
        font_size: Font size for page numbers.
        skip_first: Number of pages to skip from the beginning (e.g., cover page).
        fmt: 'decimal' (1,2,3), 'roman' (I,II,III), 'alpha' (A,B,C).
        password: PDF password if encrypted.

    Returns:
        Path to the numbered PDF.
    """
    valid_positions = {
        "bottom-center", "bottom-left", "bottom-right",
        "top-center", "top-left", "top-right",
    }
    if position not in valid_positions:
        raise ValueError(f"position must be one of: {', '.join(sorted(valid_positions))}")
    if fmt not in ("decimal", "roman", "alpha"):
        raise ValueError("fmt must be 'decimal', 'roman', or 'alpha'.")
    if start_number < 1:
        raise ValueError("start_number must be >= 1.")
    if font_size < 4 or font_size > 72:
        raise ValueError("font_size must be between 4 and 72.")

    def _to_roman(n: int) -> str:
        val = [(1000, "M"), (900, "CM"), (500, "D"), (400, "CD"),
               (100, "C"), (90, "XC"), (50, "L"), (40, "XL"),
               (10, "X"), (9, "IX"), (5, "V"), (4, "IV"), (1, "I")]
        result = ""
        for v, s in val:
            while n >= v:
                result += s
                n -= v
        return result

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        doc = fitz.open(decrypted_path)
        for i, page in enumerate(doc):
            if i < skip_first:
                continue

            page_num = start_number + (i - skip_first)
            if fmt == "roman":
                label = _to_roman(page_num)
            elif fmt == "alpha":
                label = chr(64 + page_num) if page_num <= 26 else str(page_num)
            else:
                label = str(page_num)

            rect = page.rect
            margin = 20
            y_bottom = rect.height - margin
            y_top = margin + font_size

            if "bottom" in position:
                y = y_bottom
            else:
                y = y_top

            label_width = fitz.get_text_length(label, fontname="helv", fontsize=font_size)

            if "left" in position:
                x = margin
            elif "right" in position:
                x = rect.width - margin - label_width
            else:  # center
                x = rect.width / 2 - label_width / 2

            page.insert_text(
                fitz.Point(x, y),
                label,
                fontname="helv",
                fontsize=font_size,
                color=(0, 0, 0),
            )

        doc.save(str(output_file), garbage=3, deflate=True)
        doc.close()
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return str(output_file)


# ──────────────────────────────────────────────
# Feature #61: Repair PDF
# ──────────────────────────────────────────────

def repair_pdf(input_path: str, output_dir: str) -> dict:
    """Attempt to recover/repair a corrupted or damaged PDF.

    Args:
        input_path: Path to the damaged PDF.
        output_dir: Directory to save the repaired PDF.

    Returns:
        dict with output_path and repair_status message.
    """
    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    repair_status = "unknown"
    try:
        with pikepdf.open(input_file, allow_overwriting_input=False) as pdf:
            pdf.save(output_file)
            repair_status = "success"
    except pikepdf.PdfError:
        # Try with recovery mode
        try:
            with pikepdf.open(input_file, suppress_warnings=True) as pdf:
                pdf.save(output_file)
                repair_status = "partial_recovery"
        except Exception as e:
            # Try PyMuPDF as last resort
            try:
                doc = fitz.open(str(input_file))
                doc.save(str(output_file), garbage=4, deflate=True, clean=True)
                doc.close()
                repair_status = "recovered_via_mupdf"
            except Exception as e2:
                raise RuntimeError(
                    f"Could not repair PDF: {e2}. The file may be too severely damaged."
                )

    return {"output_path": str(output_file), "repair_status": repair_status}


# ──────────────────────────────────────────────
# Feature #62: Create PDF from Scratch
# ──────────────────────────────────────────────

def create_pdf_from_text(
    output_dir: str,
    content: str,
    title: str = "Document",
    font_size: int = 12,
    page_size: str = "A4",
    margin_pt: int = 72,
) -> str:
    """Create a new PDF from plain text content.

    Args:
        output_dir: Directory to save the PDF.
        content: Text content to write into the PDF.
        title: Document title (used in filename and metadata).
        font_size: Body text font size.
        page_size: 'A4' or 'Letter'.
        margin_pt: Page margin in points.

    Returns:
        Path to the created PDF.
    """
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.pagesizes import A4, LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import ParagraphStyle

    if not content or not content.strip():
        raise ValueError("Content cannot be empty.")

    page_sizes = {"a4": A4, "letter": LETTER}
    psize = page_sizes.get(page_size.lower(), A4)

    safe_title = "".join(c for c in title if c.isalnum() or c in " _-")[:50] or "document"
    output_name = f"{safe_title.replace(' ', '_')}_{uuid.uuid4().hex[:6]}.pdf"
    output_file = Path(output_dir) / output_name

    doc = SimpleDocTemplate(
        str(output_file),
        pagesize=psize,
        leftMargin=margin_pt,
        rightMargin=margin_pt,
        topMargin=margin_pt,
        bottomMargin=margin_pt,
        title=title,
    )

    styles = getSampleStyleSheet()
    body_style = ParagraphStyle(
        "body",
        parent=styles["Normal"],
        fontSize=font_size,
        leading=font_size * 1.4,
        spaceAfter=6,
    )

    elements = []
    for paragraph in content.split("\n"):
        if paragraph.strip():
            elements.append(Paragraph(paragraph.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), body_style))
        else:
            elements.append(Spacer(1, font_size * 0.5))

    doc.build(elements)
    return str(output_file)


def create_blank_pdf(
    output_dir: str,
    num_pages: int = 1,
    page_size: str = "A4",
) -> str:
    """Create a blank PDF with the specified number of pages.

    Args:
        output_dir: Directory to save the PDF.
        num_pages: Number of blank pages (1–100).
        page_size: 'A4' or 'Letter'.

    Returns:
        Path to the created PDF.
    """
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.pagesizes import A4, LETTER

    if not 1 <= num_pages <= 100:
        raise ValueError("num_pages must be between 1 and 100.")

    page_sizes = {"a4": A4, "letter": LETTER}
    psize = page_sizes.get(page_size.lower(), A4)

    output_name = f"blank_{num_pages}pages_{uuid.uuid4().hex[:6]}.pdf"
    output_file = Path(output_dir) / output_name

    c = rl_canvas.Canvas(str(output_file), pagesize=psize)
    for _ in range(num_pages):
        c.showPage()
    c.save()

    return str(output_file)


# ──────────────────────────────────────────────
# Feature #63: Annotate / Edit PDF
# ──────────────────────────────────────────────

def annotate_pdf(
    input_path: str,
    output_dir: str,
    annotations: list,
    password: str = None,
) -> str:
    """Add annotations (highlight, underline, strikeout, note, text box, redact) to a PDF.

    Each annotation dict must contain:
        type: 'highlight' | 'underline' | 'strikeout' | 'note' | 'text' | 'redact'
        page: 1-based page number
        rect: [x0, y0, x1, y1] in PDF points
        content: text for 'note' / 'text' types (optional for others)
        color: optional [r,g,b] floats 0-1 (used for highlight color)

    Args:
        input_path: Path to input PDF.
        output_dir: Directory to save annotated PDF.
        annotations: List of annotation dicts.
        password: PDF password if encrypted.

    Returns:
        Path to the annotated PDF.
    """
    VALID_TYPES = {"highlight", "underline", "strikeout", "note", "text", "redact"}

    if not isinstance(annotations, list):
        raise ValueError("annotations must be a list.")

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        doc = fitz.open(decrypted_path)
        total_pages = len(doc)

        for ann in annotations:
            ann_type = ann.get("type", "").lower()
            if ann_type not in VALID_TYPES:
                raise ValueError(f"Unknown annotation type '{ann_type}'. Must be one of: {', '.join(sorted(VALID_TYPES))}")

            page_num = int(ann.get("page", 1))
            if page_num < 1 or page_num > total_pages:
                raise ValueError(f"Page {page_num} is out of range (document has {total_pages} pages).")

            page = doc[page_num - 1]
            raw_rect = ann.get("rect", [0, 0, 100, 20])
            rect = fitz.Rect(*raw_rect)
            content = ann.get("content", "")
            color = ann.get("color", [1, 1, 0])  # default yellow for highlights

            if ann_type == "highlight":
                page.add_highlight_annot(rect).set_colors(stroke=color)
            elif ann_type == "underline":
                page.add_underline_annot(rect)
            elif ann_type == "strikeout":
                page.add_strikeout_annot(rect)
            elif ann_type == "note":
                page.add_text_annot(rect.tl, content)
            elif ann_type == "text":
                page.insert_textbox(rect, content, fontname="helv", fontsize=11, color=(0, 0, 0))
            elif ann_type == "redact":
                # `fill=(0, 0, 0)` is load-bearing: apply_redactions() on the pinned
                # PyMuPDF<1.24 (no `graphics` param yet — added 1.23.27) strips text and
                # blanks overlapping raster images, but leaves vector content (drawn
                # rects/lines, form-field borders) under the rect fully visible with no
                # fill. The black box covers that gap at the render/extraction level even
                # though the underlying vector path stays in the content stream.
                page.add_redact_annot(rect, fill=(0, 0, 0))
                page.apply_redactions()

        doc.save(str(output_file), garbage=3, deflate=True)
        doc.close()
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return str(output_file)


# ──────────────────────────────────────────────
# Feature #64: PDF Metadata Editor
# ──────────────────────────────────────────────

def edit_pdf_metadata(
    input_path: str,
    output_dir: str,
    title: str = None,
    author: str = None,
    subject: str = None,
    keywords: str = None,
    creator: str = None,
    clear_all: bool = False,
    password: str = None,
) -> str:
    """View and edit PDF metadata (document properties).

    Args:
        input_path: Path to input PDF.
        output_dir: Directory to save the updated PDF.
        title: New document title (None = keep existing).
        author: New author name.
        subject: New subject.
        keywords: New keywords string.
        creator: New creator application name.
        clear_all: If True, remove all existing metadata before applying new values.
        password: PDF password if encrypted.

    Returns:
        Path to the updated PDF.
    """
    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)

    try:
        with pikepdf.open(decrypted_path) as pdf:
            with pdf.open_metadata(set_pikepdf_as_editor=False) as meta:
                if clear_all:
                    # Remove standard XMP metadata fields
                    for key in list(meta.keys()):
                        try:
                            del meta[key]
                        except Exception:
                            pass

                field_map = {
                    "dc:title": title,
                    "dc:creator": author,
                    "dc:description": subject,
                    "pdf:Keywords": keywords,
                    "xmp:CreatorTool": creator,
                }
                for xmp_key, value in field_map.items():
                    if value is not None:
                        meta[xmp_key] = value

            # Also update the classic docinfo dict for compatibility
            docinfo = pdf.docinfo
            if title is not None:
                docinfo["/Title"] = title
            if author is not None:
                docinfo["/Author"] = author
            if subject is not None:
                docinfo["/Subject"] = subject
            if keywords is not None:
                docinfo["/Keywords"] = keywords
            if creator is not None:
                docinfo["/Creator"] = creator

            pdf.save(output_file)
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)

    return str(output_file)


def get_pdf_metadata(input_path: str, password: str = None) -> dict:
    """Read metadata from a PDF without modifying it.

    Args:
        input_path: Path to input PDF.
        password: PDF password if encrypted.

    Returns:
        dict with metadata fields.
    """
    decrypted_path, needs_cleanup = _get_decrypted_pdf_path(input_path, password)
    try:
        with pikepdf.open(decrypted_path) as pdf:
            docinfo = pdf.docinfo
            return {
                "title": str(docinfo.get("/Title", "")),
                "author": str(docinfo.get("/Author", "")),
                "subject": str(docinfo.get("/Subject", "")),
                "keywords": str(docinfo.get("/Keywords", "")),
                "creator": str(docinfo.get("/Creator", "")),
                "producer": str(docinfo.get("/Producer", "")),
                "page_count": len(pdf.pages),
            }
    finally:
        if needs_cleanup:
            Path(decrypted_path).unlink(missing_ok=True)
