"""
PowerPoint utilities for Forge Files.

Pure-Python implementations using python-pptx + Pillow + reportlab.

NOTE: pure-Python PPT rendering is best-effort. Each slide is rendered onto a Pillow
canvas using only text, basic shape positions, and image content from the .pptx XML.
Animations, gradients, themes, SmartArt, and complex effects are NOT preserved.
For pixel-perfect rendering, install LibreOffice and shell out to `soffice`.
"""
import io
import uuid
import zipfile
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.utils import ImageReader

from scripts.utils import branded_filename, original_stem


# Render slides at 96 DPI (1 EMU = 1/914400 inch -> 1 EMU = 96/914400 px = 1/9525 px).
_EMU_PER_PX = 9525


def _emu_to_px(emu_value: Optional[int]) -> int:
    if emu_value is None:
        return 0
    return int(emu_value / _EMU_PER_PX)


def _try_font(size: int) -> ImageFont.FreeTypeFont:
    for name in ("arial.ttf", "DejaVuSans.ttf", "Helvetica.ttf"):
        try:
            return ImageFont.truetype(name, size)
        except (OSError, IOError):
            continue
    return ImageFont.load_default()


def _render_slide_to_image(slide: Slide, slide_width_emu: Optional[int], slide_height_emu: Optional[int]) -> Image.Image:
    """Best-effort raster of one slide: solid white background, then text/images at their layout positions."""
    w_px = max(1, _emu_to_px(slide_width_emu))
    h_px = max(1, _emu_to_px(slide_height_emu))
    img = Image.new("RGB", (w_px, h_px), "white")
    draw = ImageDraw.Draw(img)

    for shape in slide.shapes:
        try:
            x = _emu_to_px(shape.left or 0)
            y = _emu_to_px(shape.top or 0)
            sw = _emu_to_px(shape.width or 0)
            sh = _emu_to_px(shape.height or 0)
        except Exception:
            x = y = sw = sh = 0

        # Embedded picture — paste it at the shape rect.
        try:
            if shape.shape_type == 13 and hasattr(shape, "image"):  # MSO_SHAPE_TYPE.PICTURE = 13
                pic_bytes = shape.image.blob
                with Image.open(io.BytesIO(pic_bytes)) as pic:
                    pic = pic.convert("RGB")
                    if sw > 0 and sh > 0:
                        pic = pic.resize((sw, sh), Image.Resampling.LANCZOS)
                    img.paste(pic, (x, y))
                continue
        except Exception:
            pass

        # Text frame — render each paragraph at the shape position.
        if shape.has_text_frame:
            cur_y = y
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs) or para.text or ""
                if not text.strip():
                    cur_y += 24
                    continue

                # Pick a font size: use the first run's size if available, else 18pt.
                font_size_pt = 18
                if para.runs and para.runs[0].font.size is not None:
                    font_size_pt = max(8, int(para.runs[0].font.size.pt))
                font_px = max(10, int(font_size_pt * 1.33))  # pt -> px @ 96 DPI
                font = _try_font(font_px)

                color = (20, 20, 20)
                try:
                    if para.runs and para.runs[0].font.color and para.runs[0].font.color.rgb:
                        rgb = para.runs[0].font.color.rgb
                        color = (rgb[0], rgb[1], rgb[2])
                except Exception:
                    pass

                draw.text((x, cur_y), text, font=font, fill=color)
                cur_y += font_px + 6

    return img


def ppt_to_images_zip(input_path: str, output_dir: str, fmt: str = "png") -> dict:
    """Render every slide of a .pptx as an image and return a zip."""
    fmt = (fmt or "png").lower()
    if fmt not in ("png", "jpg", "jpeg"):
        raise ValueError("fmt must be png or jpg.")
    pil_fmt = "PNG" if fmt == "png" else "JPEG"
    file_ext = "png" if fmt == "png" else "jpg"

    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "zip")

    prs = Presentation(str(input_file))
    slide_count = len(prs.slides)

    with zipfile.ZipFile(output_file, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for i, slide in enumerate(prs.slides, start=1):
            img = _render_slide_to_image(slide, prs.slide_width, prs.slide_height)
            buf = io.BytesIO()
            if pil_fmt == "JPEG":
                img.save(buf, pil_fmt, quality=88, optimize=True)
            else:
                img.save(buf, pil_fmt, optimize=True)
            arcname = Path(f"{original_stem(input_file)}_slide_{i:03d}.{file_ext}").name
            zf.writestr(arcname, buf.getvalue())

    return {"output_path": str(output_file), "slide_count": slide_count}


def ppt_to_pdf(input_path: str, output_dir: str) -> str:
    """Render every slide as an image and assemble into a PDF."""
    input_file = Path(input_path)
    output_file = Path(output_dir) / branded_filename(input_file, "pdf")

    prs = Presentation(str(input_file))
    if not list(prs.slides):
        raise ValueError("Presentation has no slides.")

    # Use the slide dimensions in points (EMU/12700).
    page_w_pt = (prs.slide_width or 9144000) / 12700
    page_h_pt = (prs.slide_height or 6858000) / 12700

    c = rl_canvas.Canvas(str(output_file), pagesize=(page_w_pt, page_h_pt))
    for slide in prs.slides:
        img = _render_slide_to_image(slide, prs.slide_width, prs.slide_height)
        buf = io.BytesIO()
        img.save(buf, "PNG", optimize=True)
        buf.seek(0)
        c.drawImage(ImageReader(buf), 0, 0, width=page_w_pt, height=page_h_pt)
        c.showPage()
    c.save()

    return str(output_file)


def merge_pptx(input_paths: List[str], output_dir: str) -> str:
    """Merge multiple .pptx files by appending slides into the first deck.

    LIMITATION: this implementation deep-copies shape XML only — it does NOT
    copy referenced package parts (images, embedded charts, media) or remap
    their relationship IDs. As a result, slides from non-first decks that
    contain pictures/charts/media will render with broken references in the
    merged file. Text-only and shape-only decks merge cleanly. A full-fidelity
    merge would either need to walk r:id relationships and copy the underlying
    parts into the destination package, or rebuild each shape via the
    high-level python-pptx APIs (add_picture, add_chart, etc.). See PR
    discussion for context.
    """
    if not input_paths:
        raise ValueError("No input files provided for merging.")
    if len(input_paths) < 2:
        raise ValueError("Provide at least two PPTX files to merge.")

    from copy import deepcopy
    from lxml import etree

    output_file = Path(output_dir) / f"merged_{uuid.uuid4().hex[:8]}.pptx"
    base = Presentation(str(input_paths[0]))

    for path in input_paths[1:]:
        src = Presentation(str(path))
        for slide in src.slides:
            # Use the destination's first slide layout — best-effort to keep things rendering.
            blank_layout = base.slide_layouts[6] if len(base.slide_layouts) > 6 else base.slide_layouts[0]
            new_slide = base.slides.add_slide(blank_layout)

            # Strip the placeholders the layout added so we don't double-stack content.
            for shape in list(new_slide.shapes):
                sp = shape._element
                sp.getparent().remove(sp)

            # Deep-copy each shape's XML from the source slide.
            for shape in slide.shapes:
                new_slide.shapes._spTree.insert_element_before(deepcopy(shape._element), "p:extLst")

    base.save(output_file)
    return str(output_file)
