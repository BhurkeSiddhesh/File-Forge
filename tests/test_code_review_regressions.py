"""Regression tests documenting defects found in the 2026-07 full-repo code review.

Each test encodes the *correct* behavior and is marked ``xfail(strict=True)``
because the current code does not deliver it. When a fix lands, the test flips
to XPASS, pytest errors on the strict marker, and the marker must be removed —
so these tests keep the suite green today and become real guards the moment the
bugs are fixed.

Covered findings:
  * /api/workflow/execute returns 500 (unhandled ValueError) for valid-JSON but
    empty/invalid step lists — the try block only catches JSONDecodeError.
  * The MAX_UPLOAD_MB size cap lives only in save_upload(), which exactly one
    endpoint (/api/pdf/remove-password) calls; every other upload endpoint
    writes the body to disk unbounded.
  * pdf_to_pptx mutates the presentation-global slide size once per page, so in
    a mixed-page-size PDF every slide except the last is mis-scaled.
"""
import json

import pytest
from fastapi.testclient import TestClient

import main as main_module
from main import app


@pytest.fixture
def tolerant_client():
    """Client that reports 500s as responses instead of raising them."""
    return TestClient(app, raise_server_exceptions=False)


class TestWorkflowStepValidation:
    """steps='[]' passes json.loads, then `raise ValueError('No steps provided')`
    escapes the try block (which only catches json.JSONDecodeError) and becomes
    an unhandled 500 for a plain client-input mistake."""

    @pytest.mark.xfail(
        strict=True,
        reason="BUG: empty steps list raises uncaught ValueError -> 500; should be 400",
    )
    def test_empty_steps_list_is_client_error(self, tolerant_client, sample_pdf):
        with open(sample_pdf, "rb") as f:
            r = tolerant_client.post(
                "/api/workflow/execute",
                files={"file": ("sample.pdf", f, "application/pdf")},
                data={"steps": "[]"},
            )
        assert r.status_code == 400

    @pytest.mark.xfail(
        strict=True,
        reason="BUG: non-list steps JSON (e.g. '0') raises uncaught ValueError -> 500; should be 400",
    )
    def test_non_list_steps_json_is_client_error(self, tolerant_client, sample_pdf):
        with open(sample_pdf, "rb") as f:
            r = tolerant_client.post(
                "/api/workflow/execute",
                files={"file": ("sample.pdf", f, "application/pdf")},
                data={"steps": "0"},
            )
        assert r.status_code == 400


class TestUploadSizeCapCoverage:
    """MAX_UPLOAD_MB (413) and the extension allowlist are enforced by
    save_upload()/save_uploads(), which every upload endpoint now routes
    through — not just /api/pdf/remove-password. See test_input_validation.py
    for the parametrised sweep over the whole API surface."""

    def test_rotate_honors_upload_size_cap(self, tolerant_client, sample_pdf, monkeypatch):
        # With a 0 MB cap, any non-empty upload must be rejected with 413 —
        # exactly what /api/pdf/remove-password (via save_upload) already does.
        monkeypatch.setattr(main_module, "MAX_UPLOAD_MB", 0)
        with open(sample_pdf, "rb") as f:
            r = tolerant_client.post(
                "/api/pdf/rotate",
                files={"file": ("sample.pdf", f, "application/pdf")},
                data={"angle": "90"},
            )
        assert r.status_code == 413

    def test_remove_password_honors_upload_size_cap(self, tolerant_client, sample_pdf, monkeypatch):
        """Control: the one save_upload() endpoint does enforce the cap."""
        monkeypatch.setattr(main_module, "MAX_UPLOAD_MB", 0)
        with open(sample_pdf, "rb") as f:
            r = tolerant_client.post(
                "/api/pdf/remove-password",
                files={"file": ("sample.pdf", f, "application/pdf")},
                data={"password": "irrelevant"},
            )
        assert r.status_code == 413


class TestPdfToPptxMixedPageSizes:
    """FIXED (2026-08-01): a .pptx deck can carry only ONE slide size, so
    pdf_to_pptx now sizes the deck once to the bounding box of every page and
    places each page's image at its native size, centered — preserving each
    page's aspect ratio (no stretching) and clipping nothing. Previously the
    slide size was mutated per page inside the loop, so only the last page's
    size survived and every earlier slide's full-page image no longer matched
    its canvas.

    The old xfail here asserted the image should *fill* its slide; on a
    mixed-size deck that would force a stretch and distort off-size pages,
    violating the feature's "no visible raster loss" bar. The guard below now
    encodes the correct behavior: aspect preserved and nothing clipped."""

    @pytest.fixture
    def mixed_size_pdf(self, tmp_path):
        from reportlab.pdfgen import canvas

        path = tmp_path / "mixed.pdf"
        c = canvas.Canvas(str(path), pagesize=(595, 842))  # portrait A4
        c.drawString(72, 700, "portrait page")
        c.showPage()
        c.setPageSize((842, 595))  # landscape A4
        c.drawString(72, 400, "landscape page")
        c.showPage()
        c.save()
        return path

    @pytest.fixture
    def uniform_pdf(self, tmp_path):
        from reportlab.pdfgen import canvas

        path = tmp_path / "uniform.pdf"
        c = canvas.Canvas(str(path), pagesize=(595, 842))  # two portrait A4 pages
        c.drawString(72, 700, "page one")
        c.showPage()
        c.drawString(72, 700, "page two")
        c.showPage()
        c.save()
        return path

    def test_mixed_pages_preserve_aspect_and_are_not_clipped(self, mixed_size_pdf, tmp_path):
        import fitz
        from pptx import Presentation

        from scripts.pdf_utils import pdf_to_pptx

        out = pdf_to_pptx(str(mixed_size_pdf), str(tmp_path))

        # Expected per-page aspect ratios, straight from the source PDF.
        src = fitz.open(str(mixed_size_pdf))
        expected_aspect = [p.rect.width / p.rect.height for p in src]
        src.close()

        prs = Presentation(out)
        slides = list(prs.slides)
        assert len(slides) == 2
        for idx, slide in enumerate(slides):
            pictures = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
            assert pictures, "each slide should carry the rendered page image"
            pic = pictures[0]
            # Aspect ratio matches the source page (no stretch/distortion).
            assert pic.width / pic.height == pytest.approx(expected_aspect[idx], rel=1e-3)
            # Fully on-canvas: nothing clipped, positioned within the slide.
            assert pic.left >= 0 and pic.top >= 0
            assert pic.left + pic.width <= prs.slide_width
            assert pic.top + pic.height <= prs.slide_height

    def test_uniform_pages_still_fill_their_slide(self, uniform_pdf, tmp_path):
        """The common case (all pages same size) must be unchanged: every image
        fills its slide exactly at (0, 0)."""
        from pptx import Presentation

        from scripts.pdf_utils import pdf_to_pptx

        out = pdf_to_pptx(str(uniform_pdf), str(tmp_path))
        prs = Presentation(out)
        for slide in prs.slides:
            pic = [s for s in slide.shapes if s.shape_type == 13][0]
            assert (pic.left, pic.top) == (0, 0)
            assert (pic.width, pic.height) == (prs.slide_width, prs.slide_height)
