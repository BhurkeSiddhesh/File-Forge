import csv
import zipfile
from pathlib import Path

import pytest
from openpyxl import Workbook, load_workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from scripts.excel_utils import (
    csv_to_xlsx,
    excel_to_pdf,
    merge_excel_files,
    xlsx_to_csv,
)
from scripts.image_utils import (
    compress_image,
    convert_image_format,
    rotate_image,
    watermark_image,
)
from scripts.ppt_utils import merge_pptx, ppt_to_images_zip, ppt_to_pdf


# --- Fixtures ---

@pytest.fixture
def jpeg_image(tmp_path):
    p = tmp_path / "in.jpg"
    Image.new("RGB", (200, 100), color=(120, 80, 200)).save(p, "JPEG", quality=95)
    return p


@pytest.fixture
def png_image(tmp_path):
    p = tmp_path / "in.png"
    Image.new("RGBA", (160, 120), color=(255, 0, 0, 200)).save(p, "PNG")
    return p


@pytest.fixture
def csv_file(tmp_path):
    p = tmp_path / "data.csv"
    with p.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["name", "age", "city"])
        w.writerow(["alice", 30, "NYC"])
        w.writerow(["bob", 28, "SF"])
    return p


@pytest.fixture
def xlsx_file(tmp_path):
    p = tmp_path / "data.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "People"
    ws.append(["name", "age"])
    ws.append(["alice", 30])
    ws.append(["bob", 28])
    wb.save(p)
    return p


@pytest.fixture
def two_xlsx_files(tmp_path):
    paths = []
    for i in range(2):
        p = tmp_path / f"book_{i}.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = f"Sheet_{i}"
        ws.append([f"file{i}_a", f"file{i}_b"])
        ws.append([1, 2])
        wb.save(p)
        paths.append(p)
    return paths


@pytest.fixture
def pptx_file(tmp_path):
    p = tmp_path / "deck.pptx"
    prs = Presentation()
    for i in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
        title = slide.shapes.title
        title.text = f"Slide {i + 1}"
        # Add a freeform text box.
        tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        tf = tx.text_frame
        tf.text = f"Body text {i + 1}"
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(20)
    prs.save(p)
    return p


@pytest.fixture
def two_pptx_files(tmp_path):
    paths = []
    for i in range(2):
        p = tmp_path / f"deck_{i}.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Deck {i} slide"
        prs.save(p)
        paths.append(p)
    return paths


# --- Image tests ---

def test_rotate_image_90_swaps_dimensions(jpeg_image, tmp_path):
    out = Path(rotate_image(str(jpeg_image), str(tmp_path), 90))
    assert out.exists()
    with Image.open(out) as img:
        assert img.size == (100, 200)


def test_rotate_image_rejects_bad_angle(jpeg_image, tmp_path):
    with pytest.raises(ValueError):
        rotate_image(str(jpeg_image), str(tmp_path), "abc")


def test_compress_image_reduces_size(jpeg_image, tmp_path):
    # Make the image a real photo-like one so JPEG q=10 actually compresses.
    big = tmp_path / "big.jpg"
    Image.new("RGB", (800, 800), (200, 100, 50)).save(big, "JPEG", quality=95)
    result = compress_image(str(big), str(tmp_path), quality=10)
    assert Path(result["output_path"]).exists()
    assert result["compressed_size"] <= result["original_size"]


def test_compress_image_rejects_bad_quality(jpeg_image, tmp_path):
    with pytest.raises(ValueError):
        compress_image(str(jpeg_image), str(tmp_path), quality=200)


def test_convert_image_to_png(jpeg_image, tmp_path):
    out = Path(convert_image_format(str(jpeg_image), str(tmp_path), "png"))
    assert out.exists()
    assert out.suffix == ".png"
    with Image.open(out) as img:
        assert img.format == "PNG"


def test_convert_image_to_webp(png_image, tmp_path):
    out = Path(convert_image_format(str(png_image), str(tmp_path), "webp"))
    assert out.exists()
    assert out.suffix == ".webp"


def test_convert_image_rejects_bad_format(jpeg_image, tmp_path):
    with pytest.raises(ValueError):
        convert_image_format(str(jpeg_image), str(tmp_path), "tiff")


def test_watermark_image_creates_output(jpeg_image, tmp_path):
    out = Path(watermark_image(str(jpeg_image), str(tmp_path), "DRAFT", position="diagonal"))
    assert out.exists()
    with Image.open(out) as img:
        assert img.size == (200, 100)


def test_watermark_image_rejects_empty_text(jpeg_image, tmp_path):
    with pytest.raises(ValueError):
        watermark_image(str(jpeg_image), str(tmp_path), "  ")


def test_try_font_falls_back_past_missing_arial_to_a_scalable_font(monkeypatch):
    """arial.ttf is a Windows font, almost never present on the Linux deploy
    target. try_font() must keep trying scalable fonts rather than dropping
    straight to load_default()'s tiny fixed-size bitmap font, which ignores the
    requested size entirely."""
    from PIL import ImageFont
    from scripts.utils import try_font

    real_truetype = ImageFont.truetype

    def fake_truetype(name, size):
        if name == "arial.ttf":
            raise OSError("cannot open resource")
        return real_truetype(name, size)

    monkeypatch.setattr(ImageFont, "truetype", fake_truetype)
    font = try_font(150)
    assert isinstance(font, ImageFont.FreeTypeFont)
    assert font.size == 150


# --- Excel tests ---

def test_csv_to_xlsx_roundtrip(csv_file, tmp_path):
    out = Path(csv_to_xlsx(str(csv_file), str(tmp_path)))
    assert out.exists()
    wb = load_workbook(out)
    ws = wb.active
    rows = list(ws.iter_rows(values_only=True))
    assert rows[0] == ("name", "age", "city")
    assert rows[1][0] == "alice"


def test_csv_to_xlsx_rejects_multichar_delimiter(csv_file, tmp_path):
    # "||" used to slip past the old `len > 2` check and reach csv.reader,
    # which raised a raw TypeError. Now it must raise a clear ValueError.
    with pytest.raises(ValueError, match="single character"):
        csv_to_xlsx(str(csv_file), str(tmp_path), delimiter="||")


def test_csv_to_xlsx_normalizes_tab_escape(tmp_path):
    p = tmp_path / "tabbed.csv"
    p.write_text("a\tb\n1\t2\n", encoding="utf-8")
    out = Path(csv_to_xlsx(str(p), str(tmp_path), delimiter="\\t"))
    assert out.exists()
    wb = load_workbook(out)
    rows = list(wb.active.iter_rows(values_only=True))
    assert rows[0] == ("a", "b")
    assert rows[1] == ("1", "2")


def test_xlsx_to_csv_first_sheet(xlsx_file, tmp_path):
    out = Path(xlsx_to_csv(str(xlsx_file), str(tmp_path)))
    assert out.exists()
    text = out.read_text(encoding="utf-8")
    assert "name,age" in text
    assert "alice,30" in text


def test_xlsx_to_csv_named_sheet(xlsx_file, tmp_path):
    out = Path(xlsx_to_csv(str(xlsx_file), str(tmp_path), sheet="People"))
    assert out.exists()


def test_xlsx_to_csv_unknown_sheet(xlsx_file, tmp_path):
    with pytest.raises(ValueError):
        xlsx_to_csv(str(xlsx_file), str(tmp_path), sheet="Nope")


def test_excel_to_pdf_creates_pdf(xlsx_file, tmp_path):
    out = Path(excel_to_pdf(str(xlsx_file), str(tmp_path)))
    assert out.exists()
    assert out.suffix == ".pdf"
    assert out.stat().st_size > 0


# --- Office→PDF: LibreOffice-first with pure-Python fallback ---

import shutil as _shutil  # noqa: E402

_HAVE_LIBREOFFICE = bool(_shutil.which("libreoffice") or _shutil.which("soffice"))


def _pdf_text(path):
    """Extract all text from a PDF (empty string if PyMuPDF unavailable)."""
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    try:
        return "\n".join(page.get_text() for page in doc)
    finally:
        doc.close()


@pytest.fixture
def styled_xlsx_file(tmp_path):
    """A workbook exercising the fidelity the reportlab fallback drops:
    fill colors, bold fonts, a merged-cell span, and a second sheet."""
    from openpyxl.styles import Font, PatternFill

    p = tmp_path / "styled.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws["A1"] = "Region"
    ws["B1"] = "Revenue"
    ws["A1"].fill = PatternFill("solid", fgColor="FF0000")
    ws["A1"].font = Font(bold=True, color="FFFFFF")
    ws["A2"] = "North"
    ws["B2"] = 1234
    ws.merge_cells("A4:B4")
    ws["A4"] = "MERGED TOTAL"
    ws2 = wb.create_sheet("Notes")
    ws2["A1"] = "Second sheet content"
    wb.save(p)
    return p


def test_excel_to_pdf_fallback_when_libreoffice_absent(xlsx_file, tmp_path, monkeypatch):
    """When LibreOffice is unavailable the pure-Python reportlab renderer must
    still produce a valid, non-empty PDF at the branded path — the zero-regression
    guarantee for environments without LibreOffice."""
    monkeypatch.setattr("scripts.excel_utils.libreoffice_to_pdf", lambda *a, **k: None)
    out = Path(excel_to_pdf(str(xlsx_file), str(tmp_path)))
    assert out.exists()
    assert out.name == "data_forgefiles.org.pdf"
    assert out.stat().st_size > 0
    # Fallback still recovers every cell's text.
    text = _pdf_text(out)
    assert "alice" in text and "bob" in text


@pytest.mark.skipif(not _HAVE_LIBREOFFICE, reason="LibreOffice not installed")
def test_excel_to_pdf_libreoffice_preserves_content(styled_xlsx_file, tmp_path):
    """LibreOffice output keeps every sheet's cell text (incl. merged cells and a
    second sheet) and stays text-selectable (not rasterized)."""
    out = Path(excel_to_pdf(str(styled_xlsx_file), str(tmp_path)))
    assert out.exists() and out.name == "styled_forgefiles.org.pdf"
    text = _pdf_text(out)
    for probe in ("Region", "North", "MERGED TOTAL", "Second sheet content"):
        assert probe in text, f"LibreOffice output dropped {probe!r}"


def test_ppt_to_pdf_fallback_when_libreoffice_absent(pptx_file, tmp_path, monkeypatch):
    """PPT→PDF fallback (raster) must still produce a valid PDF when LibreOffice
    is unavailable — preserves the prior behavior exactly."""
    monkeypatch.setattr("scripts.ppt_utils.libreoffice_to_pdf", lambda *a, **k: None)
    out = Path(ppt_to_pdf(str(pptx_file), str(tmp_path)))
    assert out.exists()
    assert out.name == "deck_forgefiles.org.pdf"
    assert out.stat().st_size > 0


@pytest.mark.skipif(not _HAVE_LIBREOFFICE, reason="LibreOffice not installed")
def test_ppt_to_pdf_libreoffice_preserves_selectable_text(pptx_file, tmp_path):
    """LibreOffice renders slide text as real, selectable text — unlike the raster
    fallback, which bakes everything into an image with no extractable text."""
    out = Path(ppt_to_pdf(str(pptx_file), str(tmp_path)))
    assert out.exists() and out.name == "deck_forgefiles.org.pdf"
    text = _pdf_text(out)
    assert "Slide 1" in text and "Slide 2" in text


def test_libreoffice_to_pdf_returns_none_without_binary(xlsx_file, tmp_path, monkeypatch):
    """The helper degrades gracefully to None (never raises) when no LibreOffice
    binary is on PATH, so callers can fall back."""
    from scripts import utils as _utils

    monkeypatch.setattr(_utils.shutil, "which", lambda *_a, **_k: None)
    assert _utils.libreoffice_to_pdf(str(xlsx_file), str(tmp_path)) is None


def test_merge_excel_combines_sheets(two_xlsx_files, tmp_path):
    out = Path(merge_excel_files([str(p) for p in two_xlsx_files], str(tmp_path)))
    assert out.exists()
    wb = load_workbook(out)
    # Each input had one sheet -> output has 2.
    assert len(wb.sheetnames) == 2


def test_merge_excel_requires_two(xlsx_file, tmp_path):
    with pytest.raises(ValueError):
        merge_excel_files([str(xlsx_file)], str(tmp_path))


# --- PPT tests ---

def test_ppt_to_images_zip_one_per_slide(pptx_file, tmp_path):
    result = ppt_to_images_zip(str(pptx_file), str(tmp_path), fmt="png")
    out = Path(result["output_path"])
    assert out.exists()
    assert result["slide_count"] == 2
    with zipfile.ZipFile(out) as zf:
        assert len(zf.namelist()) == 2


def test_ppt_to_pdf_creates_pdf(pptx_file, tmp_path):
    out = Path(ppt_to_pdf(str(pptx_file), str(tmp_path)))
    assert out.exists()
    assert out.suffix == ".pdf"
    assert out.stat().st_size > 0


def test_merge_pptx_combines_decks(two_pptx_files, tmp_path):
    out = Path(merge_pptx([str(p) for p in two_pptx_files], str(tmp_path)))
    assert out.exists()
    prs = Presentation(str(out))
    # Original first deck has 1 slide; appending the second deck's 1 slide -> 2 total.
    assert len(prs.slides) == 2


def test_merge_pptx_requires_two(pptx_file, tmp_path):
    with pytest.raises(ValueError):
        merge_pptx([str(pptx_file)], str(tmp_path))
