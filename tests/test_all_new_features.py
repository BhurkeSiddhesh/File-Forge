"""
Comprehensive tests for all 12 new features (#53–#64).
Runs standalone without the full app (no docxcompose dependency).
"""
import json
import sys
import tempfile
from io import BytesIO
from pathlib import Path

import pikepdf
import pytest
from PIL import Image as PILImage
from reportlab.pdfgen import canvas

sys.path.insert(0, str(Path(__file__).parent.parent))

_TEST_PW = "test-pdf-pw"  # ggignore
_TEST_OWNER_PW = "test-pdf-owner-pw"  # ggignore

from scripts.pdf_utils import (
    protect_pdf,
    images_to_pdf,
    word_to_pdf,
    pdf_to_excel,
    pdf_to_pptx,
    pdf_to_epub,
    extract_text_from_pdf,
    organize_pdf,
    add_page_numbers,
    repair_pdf,
    create_pdf_from_text,
    create_blank_pdf,
    annotate_pdf,
    edit_pdf_metadata,
    get_pdf_metadata,
)


# ──────────────────────────────────────────────
# Shared Fixtures
# ──────────────────────────────────────────────


@pytest.fixture
def simple_pdf(tmp_path):
    p = tmp_path / "simple.pdf"
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "Hello world")
    c.save()
    return p


@pytest.fixture
def multi_page_pdf(tmp_path):
    p = tmp_path / "multi.pdf"
    c = canvas.Canvas(str(p))
    for i in range(1, 5):
        c.drawString(100, 750, f"Page {i} content")
        c.showPage()
    c.save()
    return p


@pytest.fixture
def locked_pdf(tmp_path, simple_pdf):
    p = tmp_path / "locked.pdf"
    pw = "secret"
    with pikepdf.open(simple_pdf) as pdf:
        pdf.save(p, encryption=pikepdf.Encryption(user=pw, owner=pw))
    return {"path": p, "password": pw}


@pytest.fixture
def sample_image(tmp_path):
    p = tmp_path / "img.png"
    img = PILImage.new("RGB", (200, 150), color=(100, 149, 237))
    img.save(p, "PNG")
    return p


@pytest.fixture
def sample_jpeg(tmp_path):
    p = tmp_path / "img.jpg"
    img = PILImage.new("RGB", (300, 200), color=(200, 100, 50))
    img.save(p, "JPEG")
    return p


@pytest.fixture
def sample_docx(tmp_path):
    """Create a minimal DOCX file for word_to_pdf testing."""
    try:
        from docx import Document
        p = tmp_path / "sample.docx"
        doc = Document()
        doc.add_paragraph("Hello, this is a test document.")
        doc.save(str(p))
        return p
    except ImportError:
        pytest.skip("python-docx not available")


@pytest.fixture
def pdf_with_table(tmp_path):
    """Create a PDF containing tabular data using PyMuPDF."""
    import fitz
    p = tmp_path / "table.pdf"
    doc = fitz.open()
    page = doc.new_page()
    # Draw a simple table structure
    shape = page.new_shape()
    # Draw grid lines
    for row in range(4):
        y = 700 - row * 40
        shape.draw_line(fitz.Point(50, y), fitz.Point(350, y))
    for col in range(4):
        x = 50 + col * 100
        shape.draw_line(fitz.Point(x, 700), fitz.Point(x, 580))
    shape.finish(color=(0, 0, 0))
    shape.commit()
    # Add text
    for row in range(3):
        for col in range(3):
            page.insert_text(
                fitz.Point(55 + col * 100, 690 - row * 40),
                f"R{row}C{col}",
                fontsize=10,
            )
    doc.save(str(p))
    doc.close()
    return p


@pytest.fixture
def borderless_table_pdf(tmp_path):
    """PDF whose table separates columns by whitespace only — no ruling lines.

    This is exactly what the default "lines" strategy cannot see; it exercises
    the text-position fallback in pdf_to_excel.
    """
    import fitz
    p = tmp_path / "borderless.pdf"
    doc = fitz.open()
    page = doc.new_page()
    rows = [
        ("Region", "Q1", "Q2", "Q3"),
        ("North", "100", "120", "140"),
        ("South", "90", "95", "99"),
        ("East", "70", "72", "80"),
    ]
    x_positions = [72, 200, 300, 400]
    y = 100
    for r in rows:
        for x, cell in zip(x_positions, r):
            page.insert_text((x, y), cell, fontsize=12)
        y += 30
    doc.save(str(p))
    doc.close()
    return p


@pytest.fixture
def prose_pdf(tmp_path):
    """A page of ordinary prose — must NOT be mistaken for a table."""
    import fitz
    p = tmp_path / "prose.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_textbox(
        fitz.Rect(72, 72, 500, 400),
        "This is an ordinary paragraph of prose text that flows across the page "
        "without any tabular structure whatsoever. It should not be mistaken for "
        "a table by the text strategy because the words do not align into columns.",
        fontsize=12,
    )
    doc.save(str(p))
    doc.close()
    return p


# ──────────────────────────────────────────────
# Feature #53: Protect PDF
# ──────────────────────────────────────────────

class TestProtectPDF:
    def test_protect_creates_file(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = protect_pdf(str(simple_pdf), str(out), user_password=_TEST_PW)
        assert Path(result).exists()
        assert "_forgefiles.org.pdf" in result

    def test_protected_pdf_requires_password(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = protect_pdf(str(simple_pdf), str(out), user_password=_TEST_PW)
        with pytest.raises(pikepdf.PasswordError):
            pikepdf.open(result)

    def test_protected_pdf_opens_with_correct_password(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = protect_pdf(str(simple_pdf), str(out), user_password=_TEST_PW)
        with pikepdf.open(result, password=_TEST_PW) as pdf:
            assert len(pdf.pages) >= 1

    def test_protect_empty_password_raises(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError, match="password"):
            protect_pdf(str(simple_pdf), str(out), user_password="")

    def test_protect_already_locked_pdf_with_password(self, locked_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = protect_pdf(
            str(locked_pdf["path"]), str(out),
            user_password=_TEST_PW,
            password=locked_pdf["password"],
        )
        assert Path(result).exists()

    def test_protect_preserves_page_count(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = protect_pdf(str(multi_page_pdf), str(out), user_password=_TEST_PW)
        with pikepdf.open(result, password=_TEST_PW) as pdf:
            assert len(pdf.pages) == 4

    def test_protect_with_owner_password(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = protect_pdf(str(simple_pdf), str(out), user_password=_TEST_PW, owner_password=_TEST_OWNER_PW)
        with pikepdf.open(result, password=_TEST_PW) as pdf:
            assert len(pdf.pages) >= 1


# ──────────────────────────────────────────────
# Feature #54: Image to PDF
# ──────────────────────────────────────────────

class TestImagesToPDF:
    def test_single_image_creates_pdf(self, sample_image, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = images_to_pdf([str(sample_image)], str(out))
        assert Path(result).exists()
        assert result.endswith(".pdf")
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) == 1

    def test_multiple_images_multi_page(self, sample_image, sample_jpeg, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = images_to_pdf([str(sample_image), str(sample_jpeg)], str(out))
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) == 2

    def test_no_images_raises(self, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError):
            images_to_pdf([], str(out))

    def test_letter_page_size(self, sample_image, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = images_to_pdf([str(sample_image)], str(out), page_size="Letter")
        assert Path(result).exists()

    def test_auto_page_size(self, sample_image, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = images_to_pdf([str(sample_image)], str(out), page_size="auto")
        assert Path(result).exists()

    # --- EXIF-orientation regression (reportlab drawImage ignores Orientation) ---

    @staticmethod
    def _exif_oriented_jpeg(path, orientation):
        """A landscape 200x100 JPEG stored with a non-identity EXIF Orientation,
        so a correct renderer must display it rotated (portrait 100x200)."""
        img = PILImage.new("RGB", (200, 100), color=(240, 240, 240))
        for y in range(100):
            for x in range(20):
                img.putpixel((x, y), (220, 20, 20))  # left stripe -> detectable
        exif = PILImage.Exif()
        exif[0x0112] = orientation
        img.save(path, "JPEG", exif=exif.tobytes(), quality=95)
        return path

    def test_exif_orientation_6_is_applied(self, tmp_path):
        """Orientation=6 (rotate 90 CW) must yield a portrait page, not sideways."""
        import fitz

        src = self._exif_oriented_jpeg(tmp_path / "rot6.jpg", 6)
        out = tmp_path / "out"
        out.mkdir()
        result = images_to_pdf([str(src)], str(out), page_size="auto", fit_mode="original")
        rect = fitz.open(result)[0].rect
        # Stored pixels are 200x100 landscape; with orientation applied the page
        # must be portrait (~100x200). Before the fix reportlab drew it 200x100.
        assert round(rect.width) == 100 and round(rect.height) == 200

    def test_exif_orientation_8_is_applied(self, tmp_path):
        """Orientation=8 (rotate 90 CCW) must also be baked in."""
        import fitz

        src = self._exif_oriented_jpeg(tmp_path / "rot8.jpg", 8)
        out = tmp_path / "out"
        out.mkdir()
        result = images_to_pdf([str(src)], str(out), page_size="auto", fit_mode="original")
        rect = fitz.open(result)[0].rect
        assert round(rect.width) == 100 and round(rect.height) == 200

    def test_upright_image_uses_original_untouched(self, tmp_path):
        """No-orientation inputs take the fast path unchanged: the original file is
        drawn directly (no re-encode) and no temp copy is created or left behind."""
        from scripts.pdf_utils import _oriented_for_pdf

        src = tmp_path / "plain.jpg"
        PILImage.new("RGB", (200, 100), color=(0, 0, 200)).save(src, "JPEG", quality=95)
        tmp_files = []
        draw_path, w, h = _oriented_for_pdf(str(src), str(tmp_path), tmp_files)
        assert draw_path == str(src)
        assert (w, h) == (200, 100)
        assert tmp_files == []

    def test_no_oriented_temp_files_left_behind(self, tmp_path):
        """The orientation temp copy is cleaned up after the PDF is written."""
        src = self._exif_oriented_jpeg(tmp_path / "rot6.jpg", 6)
        out = tmp_path / "out"
        out.mkdir()
        images_to_pdf([str(src)], str(out))
        assert not list(out.glob("_oriented_*"))

    def test_fit_mode_original(self, sample_image, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = images_to_pdf([str(sample_image)], str(out), fit_mode="original")
        assert Path(result).exists()

    def test_three_images(self, sample_image, sample_jpeg, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        # Use same image three times
        result = images_to_pdf([str(sample_image), str(sample_jpeg), str(sample_image)], str(out))
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) == 3


# ──────────────────────────────────────────────
# Feature #55: Word to PDF
# ──────────────────────────────────────────────

class TestWordToPDF:
    def test_docx_converts_to_pdf(self, sample_docx, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = word_to_pdf(str(sample_docx), str(out))
        assert Path(result).exists()
        assert result.endswith(".pdf")
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) >= 1

    def test_invalid_extension_raises(self, tmp_path):
        fake = tmp_path / "file.xyz"
        fake.write_text("not a doc")
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError, match="docx"):
            word_to_pdf(str(fake), str(out))

    def test_output_is_valid_pdf(self, sample_docx, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = word_to_pdf(str(sample_docx), str(out))
        # Verify it's a valid PDF by checking header
        with open(result, "rb") as f:
            header = f.read(5)
        assert header == b"%PDF-"

    def test_fallback_when_libreoffice_absent(self, sample_docx, tmp_path, monkeypatch):
        """Pure-Python fallback must still produce a valid, non-empty PDF at the
        branded path when LibreOffice is unavailable — zero-regression guarantee,
        same pattern already proven for excel_to_pdf/ppt_to_pdf."""
        monkeypatch.setattr("scripts.pdf_utils.libreoffice_to_pdf", lambda *a, **k: None)
        out = Path(word_to_pdf(str(sample_docx), str(tmp_path)))
        assert out.exists() and out.stat().st_size > 0
        with open(out, "rb") as f:
            assert f.read(5) == b"%PDF-"

    def test_non_docx_fallback_unsupported_when_libreoffice_absent(self, tmp_path, monkeypatch):
        """.doc/.odt/.rtf have no pure-Python fallback — must still raise clearly
        (unaffected by the LibreOffice call-path consolidation)."""
        monkeypatch.setattr("scripts.pdf_utils.libreoffice_to_pdf", lambda *a, **k: None)
        fake = tmp_path / "file.odt"
        fake.write_text("not a real odt")
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(RuntimeError, match="Pure-Python fallback only supports .docx"):
            word_to_pdf(str(fake), str(out))

    def test_uses_shared_libreoffice_helper_with_isolated_profile(self, sample_docx, tmp_path, monkeypatch):
        """word_to_pdf must go through the shared scripts.utils.libreoffice_to_pdf
        helper (isolated per-call -env:UserInstallation profile dir) rather than
        shelling out directly, so it doesn't contend on the shared LibreOffice
        profile lock with a concurrent excel_to_pdf/ppt_to_pdf conversion."""
        calls = []

        def _fake_libreoffice_to_pdf(input_path, output_dir, timeout=120):
            calls.append((input_path, output_dir))
            produced = Path(output_dir) / f"{Path(input_path).stem}.pdf"
            produced.write_bytes(b"%PDF-1.4\n%%EOF")
            return produced

        monkeypatch.setattr("scripts.pdf_utils.libreoffice_to_pdf", _fake_libreoffice_to_pdf)
        out = tmp_path / "out"
        out.mkdir()
        result = Path(word_to_pdf(str(sample_docx), str(out)))
        assert calls == [(str(sample_docx), str(out))]
        assert result.exists() and result.name.endswith("_forgefiles.org.pdf")


# ──────────────────────────────────────────────
# Feature #56: PDF to Excel
# ──────────────────────────────────────────────

class TestPDFToExcel:
    def test_creates_excel_file(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_excel(str(simple_pdf), str(out))
        assert Path(result["output_path"]).exists()
        assert result["output_path"].endswith(".xlsx")

    def test_returns_tables_found_count(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_excel(str(simple_pdf), str(out))
        assert "tables_found" in result
        assert isinstance(result["tables_found"], int)

    def test_pdf_with_table_extracts_data(self, pdf_with_table, tmp_path):
        import openpyxl
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_excel(str(pdf_with_table), str(out))
        wb = openpyxl.load_workbook(result["output_path"])
        assert len(wb.sheetnames) >= 1

    def test_multi_page_pdf_processes_all(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_excel(str(multi_page_pdf), str(out))
        assert Path(result["output_path"]).exists()

    def test_locked_pdf_with_password(self, locked_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_excel(str(locked_pdf["path"]), str(out), password=locked_pdf["password"])
        assert Path(result["output_path"]).exists()

    def test_no_tables_fallback_text(self, simple_pdf, tmp_path):
        import openpyxl
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_excel(str(simple_pdf), str(out))
        wb = openpyxl.load_workbook(result["output_path"])
        # Should have at least one sheet (text fallback or table sheet)
        assert len(wb.sheetnames) >= 1

    # ── Borderless-table recovery (text-position fallback) ──────────────

    def test_borderless_table_recovered(self, borderless_table_pdf, tmp_path):
        """Regression: whitespace-aligned tables were silently dropped by the
        default 'lines' strategy; the text-position fallback must now recover
        every cell."""
        import openpyxl
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_excel(str(borderless_table_pdf), str(out))
        assert result["tables_found"] >= 1, "borderless table was dropped"
        wb = openpyxl.load_workbook(result["output_path"])
        # A real table sheet (P#_T#), not the raw-text dump.
        assert any(name.startswith("P") and "_T" in name for name in wb.sheetnames)
        # Every source cell must survive into the workbook.
        cells = set()
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for c in row:
                    if c is not None and str(c).strip():
                        cells.add(str(c).strip())
        for expected in ("Region", "Q1", "North", "100", "South", "East", "80"):
            assert expected in cells, f"missing cell {expected!r}"

    def test_prose_not_treated_as_table(self, prose_pdf, tmp_path):
        """Guard: a page of prose must not become a spurious table — it should
        still fall through to the raw-text fallback sheet."""
        import openpyxl
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_excel(str(prose_pdf), str(out))
        assert result["tables_found"] == 0
        wb = openpyxl.load_workbook(result["output_path"])
        assert "Text Content" in wb.sheetnames

    def test_bordered_table_unchanged(self, pdf_with_table, tmp_path):
        """Regression: the ruled-table path must be untouched — the text
        fallback only runs when 'lines' finds nothing, so a bordered table is
        still extracted by the 'lines' strategy exactly as before."""
        import openpyxl
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_excel(str(pdf_with_table), str(out))
        assert result["tables_found"] >= 1
        wb = openpyxl.load_workbook(result["output_path"])
        assert any(name.startswith("P") and "_T" in name for name in wb.sheetnames)


# ──────────────────────────────────────────────
# Feature #57: PDF to PowerPoint
# ──────────────────────────────────────────────

class TestPDFToPPTX:
    def test_creates_pptx_file(self, simple_pdf, tmp_path):
        from pptx import Presentation
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_pptx(str(simple_pdf), str(out))
        assert Path(result).exists()
        assert result.endswith(".pptx")

    def test_slide_count_matches_pages(self, multi_page_pdf, tmp_path):
        from pptx import Presentation
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_pptx(str(multi_page_pdf), str(out))
        prs = Presentation(result)
        assert len(prs.slides) == 4

    def test_single_page_pdf(self, simple_pdf, tmp_path):
        from pptx import Presentation
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_pptx(str(simple_pdf), str(out))
        prs = Presentation(result)
        assert len(prs.slides) == 1

    def test_invalid_dpi_raises(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError, match="DPI"):
            pdf_to_pptx(str(simple_pdf), str(out), dpi=500)

    def test_with_locked_pdf(self, locked_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_pptx(str(locked_pdf["path"]), str(out), password=locked_pdf["password"])
        assert Path(result).exists()


# ──────────────────────────────────────────────
# Feature #60: PDF to EPUB
# ──────────────────────────────────────────────

class TestPDFToEpub:
    def test_creates_epub_file(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_epub(str(simple_pdf), str(out))
        assert Path(result).exists()
        assert result.endswith(".epub")
        assert Path(result).stem == f"{simple_pdf.stem}_forgefiles.org"

    def test_chapter_count_matches_pages(self, multi_page_pdf, tmp_path):
        from ebooklib import epub
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_epub(str(multi_page_pdf), str(out))
        book = epub.read_epub(result)
        chapters = [item for item in book.get_items_of_type(9) if item.file_name.startswith("page_")]
        assert len(chapters) == 4

    def test_extracted_text_appears_in_a_chapter(self, text_rich_pdf, tmp_path):
        from ebooklib import epub
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_epub(str(text_rich_pdf), str(out))
        book = epub.read_epub(result)
        all_content = b"".join(item.get_content() for item in book.get_items_of_type(9))
        assert b"Jordan Rivera" in all_content

    def test_with_locked_pdf(self, locked_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_epub(str(locked_pdf["path"]), str(out), password=locked_pdf["password"])
        assert Path(result).exists()

    def test_locked_pdf_without_password_raises(self, locked_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError):
            pdf_to_epub(str(locked_pdf["path"]), str(out))

    def test_scanned_pdf_falls_back_to_placeholder_without_ocr(self, scanned_like_pdf, tmp_path, monkeypatch):
        """When AI/OCR is unavailable, a scanned PDF still degrades gracefully
        to a placeholder chapter instead of raising."""
        import scripts.ocr_engine as ocr_engine
        monkeypatch.setattr(ocr_engine, "get_ocr_engine", lambda *a, **k: None)

        from ebooklib import epub
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_epub(str(scanned_like_pdf), str(out))
        book = epub.read_epub(result)
        all_content = b"".join(item.get_content() for item in book.get_items_of_type(9))
        assert b"No text found" in all_content

    def test_preserves_bold_and_italic(self, tmp_path):
        import fitz
        from ebooklib import epub

        pdf_path = tmp_path / "styled.pdf"
        doc = fitz.open()
        page = doc.new_page()
        bold_font = fitz.Font("hebo")  # Helvetica-Bold
        italic_font = fitz.Font("heit")  # Helvetica-Italic
        page.insert_font(fontname="F0", fontbuffer=bold_font.buffer)
        page.insert_font(fontname="F1", fontbuffer=italic_font.buffer)
        page.insert_text((72, 100), "Bold Heading", fontsize=14, fontname="F0")
        page.insert_text((72, 130), "Italic emphasis", fontsize=12, fontname="F1")
        page.insert_text((72, 160), "Plain text", fontsize=12)
        doc.save(str(pdf_path))
        doc.close()

        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_epub(str(pdf_path), str(out))
        book = epub.read_epub(result)
        content = b"".join(item.get_content() for item in book.get_items_of_type(9))

        assert b"<b>Bold Heading</b>" in content
        assert b"<i>Italic emphasis</i>" in content
        assert b"<p>Plain text</p>" in content

    def test_preserves_hyperlinks(self, tmp_path):
        import fitz
        from ebooklib import epub

        pdf_path = tmp_path / "linked.pdf"
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 100), "Visit our site: Example Site", fontsize=12)
        link_rect = page.search_for("Example Site")[0]
        page.insert_link({"kind": fitz.LINK_URI, "from": link_rect, "uri": "https://example.com"})
        doc.save(str(pdf_path))
        doc.close()

        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_epub(str(pdf_path), str(out))
        book = epub.read_epub(result)
        content = b"".join(item.get_content() for item in book.get_items_of_type(9))

        assert b'<a href="https://example.com">Example Site</a>' in content
        # Only the linked substring is wrapped — the rest of the line stays plain.
        assert b"Visit our site: <a" in content

    def test_generates_cover_thumbnail(self, multi_page_pdf, tmp_path):
        from ebooklib import epub
        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_epub(str(multi_page_pdf), str(out))
        book = epub.read_epub(result)
        # Verify cover item exists in EPUB
        cover_items = [item for item in book.get_items() if "cover" in item.get_name().lower()]
        assert len(cover_items) > 0

    def test_extracts_and_embeds_images(self, tmp_path):
        import fitz
        import ebooklib
        from ebooklib import epub
        from PIL import Image
        import io

        pdf_path = tmp_path / "doc_with_image.pdf"
        doc = fitz.open()
        page = doc.new_page(width=300, height=300)
        page.insert_text((50, 50), "Document with illustration", fontsize=12)

        # Create a 100x100 test image and insert into page via byte stream
        img = Image.new("RGB", (100, 100), color="blue")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        page.insert_image(fitz.Rect(50, 80, 150, 180), stream=buf.getvalue())

        doc.save(str(pdf_path))
        doc.close()

        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_epub(str(pdf_path), str(out))
        book = epub.read_epub(result)

        # Verify image item is saved in EPUB manifest
        image_items = [item for item in book.get_items_of_type(ebooklib.ITEM_IMAGE)]
        # Filter out cover.jpg to check extracted page images
        page_images = [img for img in image_items if "cover" not in img.get_name().lower()]
        assert len(page_images) >= 1

        content = b"".join(item.get_content() for item in book.get_items_of_type(9))
        assert b'<img src="images/page_1_img_' in content
        assert b"Document with illustration" in content

    def test_infers_h2_and_h3_headings(self, tmp_path):
        import fitz
        from ebooklib import epub

        pdf_path = tmp_path / "headings.pdf"
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 100), "Major Section Title", fontsize=20)
        page.insert_text((72, 140), "Subsection Topic", fontsize=15)
        page.insert_text((72, 180), "This is regular body text spanning the paragraph.", fontsize=11)
        page.insert_text((72, 210), "Another regular body sentence.", fontsize=11)
        doc.save(str(pdf_path))
        doc.close()

        out = tmp_path / "out"
        out.mkdir()
        result = pdf_to_epub(str(pdf_path), str(out))
        book = epub.read_epub(result)
        content = b"".join(item.get_content() for item in book.get_items_of_type(9))

        assert b"<h2>Major Section Title</h2>" in content
        assert b"<h3>Subsection Topic</h3>" in content
        assert b"<p>This is regular body text spanning the paragraph.</p>" in content


# ──────────────────────────────────────────────
# Feature #58: Extract Text from PDF
# ──────────────────────────────────────────────

class TestExtractText:
    def test_creates_txt_file(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = extract_text_from_pdf(str(simple_pdf), str(out))
        assert Path(result["output_path"]).exists()
        assert result["output_path"].endswith(".txt")

    def test_extracts_text_content(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = extract_text_from_pdf(str(simple_pdf), str(out))
        text = Path(result["output_path"]).read_text(encoding="utf-8")
        assert "Hello world" in text

    def test_returns_page_count(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = extract_text_from_pdf(str(multi_page_pdf), str(out))
        assert result["page_count"] == 4

    def test_preserve_layout_mode(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = extract_text_from_pdf(str(simple_pdf), str(out), preserve_layout=True)
        assert Path(result["output_path"]).exists()

    def test_page_headers_in_output(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = extract_text_from_pdf(str(multi_page_pdf), str(out))
        text = Path(result["output_path"]).read_text(encoding="utf-8")
        assert "--- Page 1 ---" in text

    def test_locked_pdf_with_password(self, locked_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = extract_text_from_pdf(str(locked_pdf["path"]), str(out), password=locked_pdf["password"])
        assert Path(result["output_path"]).exists()

    def test_output_filename_contains_text(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = extract_text_from_pdf(str(simple_pdf), str(out))
        assert "_forgefiles.org.txt" in result["output_path"]


# ──────────────────────────────────────────────
# Feature #59: Organize PDF
# ──────────────────────────────────────────────

class TestOrganizePDF:
    def test_reorder_pages(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = organize_pdf(str(multi_page_pdf), str(out), page_order=[4, 3, 2, 1])
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) == 4

    def test_delete_pages(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = organize_pdf(str(multi_page_pdf), str(out), page_order=[1, 3])
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) == 2

    def test_duplicate_pages(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = organize_pdf(str(multi_page_pdf), str(out), page_order=[1, 1, 2])
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) == 3

    def test_empty_order_raises(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError):
            organize_pdf(str(multi_page_pdf), str(out), page_order=[])

    def test_out_of_range_page_raises(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError):
            organize_pdf(str(multi_page_pdf), str(out), page_order=[1, 99])

    def test_output_filename(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = organize_pdf(str(multi_page_pdf), str(out), page_order=[1, 2])
        assert "_forgefiles.org.pdf" in result

    def test_single_page_extract(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = organize_pdf(str(multi_page_pdf), str(out), page_order=[2])
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) == 1

    def test_locked_pdf_with_password(self, locked_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = organize_pdf(str(locked_pdf["path"]), str(out), page_order=[1], password=locked_pdf["password"])
        assert Path(result).exists()


# ──────────────────────────────────────────────
# Feature #60: Add Page Numbers
# ──────────────────────────────────────────────

class TestAddPageNumbers:
    def test_creates_numbered_pdf(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = add_page_numbers(str(multi_page_pdf), str(out))
        assert Path(result).exists()
        assert "_forgefiles.org.pdf" in result

    def test_page_count_preserved(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = add_page_numbers(str(multi_page_pdf), str(out))
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) == 4

    def test_all_positions(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        positions = ["bottom-center", "bottom-left", "bottom-right",
                     "top-center", "top-left", "top-right"]
        for pos in positions:
            result = add_page_numbers(str(simple_pdf), str(out), position=pos)
            assert Path(result).exists(), f"Failed for position: {pos}"

    def test_invalid_position_raises(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError, match="position"):
            add_page_numbers(str(simple_pdf), str(out), position="middle")

    def test_roman_format(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = add_page_numbers(str(multi_page_pdf), str(out), fmt="roman")
        assert Path(result).exists()

    def test_alpha_format(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = add_page_numbers(str(multi_page_pdf), str(out), fmt="alpha")
        assert Path(result).exists()

    def test_skip_first_pages(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = add_page_numbers(str(multi_page_pdf), str(out), skip_first=1)
        assert Path(result).exists()

    def test_custom_start_number(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = add_page_numbers(str(multi_page_pdf), str(out), start_number=5)
        assert Path(result).exists()

    def test_invalid_format_raises(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError, match="fmt"):
            add_page_numbers(str(simple_pdf), str(out), fmt="unknown")

    def test_locked_pdf_with_password(self, locked_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = add_page_numbers(str(locked_pdf["path"]), str(out), password=locked_pdf["password"])
        assert Path(result).exists()

    def test_right_position_uses_exact_text_width(self, multi_page_pdf, tmp_path):
        """Right-aligned labels must sit flush with the margin using real glyph
        widths (fitz.get_text_length), not a per-character heuristic that
        drifts further off with every extra digit/font-size point."""
        import fitz

        out = tmp_path / "out"
        out.mkdir()
        margin = 20
        font_size = 12
        # start_number chosen so labels span 1-4 digits across the 4-page fixture.
        result = add_page_numbers(
            str(multi_page_pdf), str(out), position="bottom-right",
            start_number=999, font_size=font_size,
        )
        doc = fitz.open(result)
        try:
            for i, page in enumerate(doc):
                label = str(999 + i)
                expected_width = fitz.get_text_length(label, fontname="helv", fontsize=font_size)
                expected_x = page.rect.width - margin - expected_width
                spans = [
                    span for block in page.get_text("dict")["blocks"]
                    for line in block.get("lines", [])
                    for span in line["spans"]
                    if span["text"].strip() == label
                ]
                assert spans, f"page {i}: label {label!r} not found"
                actual_x = spans[0]["bbox"][0]
                assert abs(actual_x - expected_x) < 1.0, (
                    f"page {i}: label {label!r} x={actual_x} expected~={expected_x}"
                )
        finally:
            doc.close()

    def test_center_position_is_symmetric(self, simple_pdf, tmp_path):
        """Center-aligned label's left/right whitespace margins must match
        within a point, using the real text width rather than a heuristic."""
        import fitz

        out = tmp_path / "out"
        out.mkdir()
        result = add_page_numbers(str(simple_pdf), str(out), position="bottom-center")
        doc = fitz.open(result)
        try:
            page = doc[0]
            spans = [
                span for block in page.get_text("dict")["blocks"]
                for line in block.get("lines", [])
                for span in line["spans"]
                if span["text"].strip() == "1"
            ]
            assert spans
            bbox = spans[0]["bbox"]
            left_gap = bbox[0]
            right_gap = page.rect.width - bbox[2]
            assert abs(left_gap - right_gap) < 1.0
        finally:
            doc.close()


# ──────────────────────────────────────────────
# Feature #61: Repair PDF
# ──────────────────────────────────────────────

class TestRepairPDF:
    def test_repair_valid_pdf(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = repair_pdf(str(simple_pdf), str(out))
        assert Path(result["output_path"]).exists()
        assert result["repair_status"] in ("success", "partial_recovery", "recovered_via_mupdf")

    def test_repair_output_is_valid_pdf(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = repair_pdf(str(simple_pdf), str(out))
        with pikepdf.open(result["output_path"]) as pdf:
            assert len(pdf.pages) >= 1

    def test_repair_output_filename(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = repair_pdf(str(simple_pdf), str(out))
        assert "_forgefiles.org.pdf" in result["output_path"]

    def test_repair_preserves_pages(self, multi_page_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = repair_pdf(str(multi_page_pdf), str(out))
        with pikepdf.open(result["output_path"]) as pdf:
            assert len(pdf.pages) == 4

    def test_repair_truncated_pdf(self, tmp_path):
        """Test repair on a partially truncated PDF."""
        out = tmp_path / "out"
        out.mkdir()

        # Create a valid PDF, then truncate it
        good_pdf = tmp_path / "good.pdf"
        c = canvas.Canvas(str(good_pdf))
        c.drawString(100, 750, "test")
        c.save()
        good_bytes = good_pdf.read_bytes()

        bad_pdf = tmp_path / "bad.pdf"
        bad_pdf.write_bytes(good_bytes[: len(good_bytes) // 2])

        try:
            result = repair_pdf(str(bad_pdf), str(out))
            assert "output_path" in result
        except RuntimeError:
            # Expected for severely corrupted files
            pass


# ──────────────────────────────────────────────
# Feature #62: Create PDF from Scratch
# ──────────────────────────────────────────────

class TestCreatePDF:
    def test_create_from_text(self, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = create_pdf_from_text(str(out), "Hello world\nSecond line")
        assert Path(result).exists()
        assert result.endswith(".pdf")
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) >= 1

    def test_empty_content_raises(self, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError, match="empty"):
            create_pdf_from_text(str(out), "   ")

    def test_custom_title(self, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = create_pdf_from_text(str(out), "Content", title="My Report")
        assert "My_Report" in result

    def test_letter_page_size(self, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = create_pdf_from_text(str(out), "Content", page_size="Letter")
        assert Path(result).exists()

    def test_large_text_creates_multiple_pages(self, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        content = "\n".join([f"Line {i}: " + "A" * 80 for i in range(200)])
        result = create_pdf_from_text(str(out), content)
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) >= 2

    def test_create_blank_pdf(self, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = create_blank_pdf(str(out), num_pages=3)
        assert Path(result).exists()
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) == 3

    def test_create_blank_single_page(self, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = create_blank_pdf(str(out), num_pages=1)
        with pikepdf.open(result) as pdf:
            assert len(pdf.pages) == 1

    def test_blank_pdf_invalid_pages_raises(self, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError):
            create_blank_pdf(str(out), num_pages=0)
        with pytest.raises(ValueError):
            create_blank_pdf(str(out), num_pages=101)

    def test_blank_pdf_letter_size(self, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = create_blank_pdf(str(out), num_pages=2, page_size="Letter")
        assert Path(result).exists()

    def test_special_chars_in_content(self, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = create_pdf_from_text(str(out), "Hello <world> & 'more'")
        assert Path(result).exists()


# ──────────────────────────────────────────────
# Feature #63: Annotate PDF
# ──────────────────────────────────────────────

class TestAnnotatePDF:
    def _make_annot(self, ann_type, page=1, rect=None, content="test"):
        return {
            "type": ann_type,
            "page": page,
            "rect": rect or [50, 700, 300, 730],
            "content": content,
        }

    def test_highlight_annotation(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        anns = [self._make_annot("highlight")]
        result = annotate_pdf(str(simple_pdf), str(out), anns)
        assert Path(result).exists()
        assert "_forgefiles.org.pdf" in result

    def test_underline_annotation(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = annotate_pdf(str(simple_pdf), str(out), [self._make_annot("underline")])
        assert Path(result).exists()

    def test_strikeout_annotation(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = annotate_pdf(str(simple_pdf), str(out), [self._make_annot("strikeout")])
        assert Path(result).exists()

    def test_note_annotation(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = annotate_pdf(str(simple_pdf), str(out), [self._make_annot("note", content="Review this")])
        assert Path(result).exists()

    def test_text_box_annotation(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = annotate_pdf(str(simple_pdf), str(out), [self._make_annot("text", content="Added text")])
        assert Path(result).exists()

    def test_redact_annotation(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = annotate_pdf(str(simple_pdf), str(out), [self._make_annot("redact")])
        assert Path(result).exists()

    def test_redact_removes_text_and_covers_vector_graphics(self, tmp_path):
        """A redaction must not just remove text — it must also visually cover any
        vector graphics (drawn rects/lines) under the box. apply_redactions() on the
        pinned PyMuPDF<1.24 has no `graphics` param (added 1.23.27) and leaves vector
        paths in the content stream untouched; without an explicit fill the redacted
        area was previously left blank, letting anything drawn there show straight
        through post-redaction."""
        import fitz

        src = tmp_path / "vector.pdf"
        doc = fitz.open()
        page = doc.new_page(width=400, height=400)
        page.insert_text((50, 55), "SECRET-TEXT-12345", fontsize=14)
        page.draw_rect(fitz.Rect(40, 30, 250, 70), color=(1, 0, 0), fill=(1, 0, 0))
        doc.save(str(src))
        doc.close()

        out = tmp_path / "out"
        out.mkdir()
        ann = self._make_annot("redact", rect=[35, 25, 260, 75])
        result = annotate_pdf(str(src), str(out), [ann])

        doc2 = fitz.open(result)
        page2 = doc2[0]
        assert "SECRET-TEXT-12345" not in page2.get_text()
        pix = page2.get_pixmap()
        # Center and edge of the redacted rect must render black, not the red
        # vector rect that was drawn underneath.
        assert pix.pixel(140, 50) == (0, 0, 0)
        assert pix.pixel(40, 30) == (0, 0, 0)
        doc2.close()

    def test_multiple_annotations(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        anns = [
            self._make_annot("highlight", rect=[50, 700, 200, 720]),
            self._make_annot("underline", rect=[50, 670, 200, 690]),
            self._make_annot("note", content="Note here"),
        ]
        result = annotate_pdf(str(simple_pdf), str(out), anns)
        assert Path(result).exists()

    def test_invalid_type_raises(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError, match="annotation type"):
            annotate_pdf(str(simple_pdf), str(out), [self._make_annot("magic")])

    def test_out_of_range_page_raises(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        with pytest.raises(ValueError, match="out of range"):
            annotate_pdf(str(simple_pdf), str(out), [self._make_annot("highlight", page=99)])

    def test_empty_annotations_list(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = annotate_pdf(str(simple_pdf), str(out), [])
        assert Path(result).exists()

    def test_locked_pdf_with_password(self, locked_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        anns = [self._make_annot("highlight")]
        result = annotate_pdf(str(locked_pdf["path"]), str(out), anns, password=locked_pdf["password"])
        assert Path(result).exists()

    def test_custom_highlight_color(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        ann = {"type": "highlight", "page": 1, "rect": [50, 700, 200, 720], "color": [0, 1, 0]}
        result = annotate_pdf(str(simple_pdf), str(out), [ann])
        assert Path(result).exists()


# ──────────────────────────────────────────────
# Feature #64: PDF Metadata Editor
# ──────────────────────────────────────────────

class TestPDFMetadata:
    def test_edit_title(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = edit_pdf_metadata(str(simple_pdf), str(out), title="My Title")
        assert Path(result).exists()

    def test_edit_author(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = edit_pdf_metadata(str(simple_pdf), str(out), author="John Doe")
        with pikepdf.open(result) as pdf:
            assert str(pdf.docinfo.get("/Author", "")) == "John Doe"

    def test_edit_multiple_fields(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = edit_pdf_metadata(
            str(simple_pdf), str(out),
            title="Test", author="Author", subject="Subject", keywords="a,b,c"
        )
        with pikepdf.open(result) as pdf:
            assert str(pdf.docinfo.get("/Title", "")) == "Test"
            assert str(pdf.docinfo.get("/Author", "")) == "Author"

    def test_read_metadata(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        # Write metadata first
        edited = edit_pdf_metadata(str(simple_pdf), str(out), title="ReadTest", author="TestAuthor")
        meta = get_pdf_metadata(edited)
        assert meta["title"] == "ReadTest"
        assert meta["author"] == "TestAuthor"
        assert "page_count" in meta
        assert meta["page_count"] >= 1

    def test_clear_all_metadata(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        out2 = tmp_path / "out2"
        out2.mkdir()
        # First set some metadata
        edited = edit_pdf_metadata(str(simple_pdf), str(out), title="ToRemove", author="ToRemove")
        # Then clear all (separate dir — same original stem would otherwise
        # re-derive the same branded name `edited` was just written to)
        cleared = edit_pdf_metadata(edited, str(out2), clear_all=True)
        assert Path(cleared).exists()

    def test_none_fields_keep_existing(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = edit_pdf_metadata(str(simple_pdf), str(out), title="Keep", author=None)
        with pikepdf.open(result) as pdf:
            assert str(pdf.docinfo.get("/Title", "")) == "Keep"

    def test_locked_pdf_with_password(self, locked_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = edit_pdf_metadata(
            str(locked_pdf["path"]), str(out),
            title="Locked Title",
            password=locked_pdf["password"]
        )
        assert Path(result).exists()

    def test_output_filename(self, simple_pdf, tmp_path):
        out = tmp_path / "out"
        out.mkdir()
        result = edit_pdf_metadata(str(simple_pdf), str(out), title="Test")
        assert "_forgefiles.org.pdf" in result

    def test_get_metadata_page_count(self, multi_page_pdf, tmp_path):
        meta = get_pdf_metadata(str(multi_page_pdf))
        assert meta["page_count"] == 4

    def test_get_metadata_returns_dict(self, simple_pdf, tmp_path):
        meta = get_pdf_metadata(str(simple_pdf))
        expected_keys = {"title", "author", "subject", "keywords", "creator", "producer", "page_count"}
        assert expected_keys.issubset(set(meta.keys()))
