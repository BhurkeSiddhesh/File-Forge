"""
Tests for Word -> PowerPoint conversion (scripts.pdf_utils.word_to_pptx), the
standalone /api/word/to-pptx endpoint, and the workflow-builder step.

word_to_pptx renders the DOCX to PDF first (word_to_pdf), then rasterizes each
page as a slide image (pdf_to_pptx) -- same non-editable, image-slide approach
pdf_to_pptx already uses for PDF input.
"""
import json
from pathlib import Path
from unittest.mock import patch

from conftest import result_path

import pytest
from pptx import Presentation

from scripts.pdf_utils import word_to_pptx


def _make_docx(tmp_path, name="report.docx", paragraphs=None):
    from docx import Document
    doc = Document()
    for text in paragraphs or ["Word to PowerPoint test content."]:
        doc.add_paragraph(text)
    p = tmp_path / name
    doc.save(str(p))
    return p


def test_word_to_pptx_creates_valid_pptx(tmp_path):
    out = tmp_path / "out"
    out.mkdir()
    docx_path = _make_docx(tmp_path)

    output_path = word_to_pptx(str(docx_path), str(out))

    assert Path(output_path).exists()
    prs = Presentation(output_path)
    assert len(prs.slides) >= 1


def test_word_to_pptx_output_name_is_branded(tmp_path):
    out = tmp_path / "out"
    out.mkdir()
    docx_path = _make_docx(tmp_path, name="Quarterly Report.docx")

    output_path = word_to_pptx(str(docx_path), str(out))

    assert Path(output_path).name == "Quarterly Report_forgefiles.org.pptx"


def test_word_to_pptx_does_not_leave_intermediate_pdf(tmp_path):
    """The temp PDF used to bridge word_to_pdf -> pdf_to_pptx must not survive."""
    out = tmp_path / "out"
    out.mkdir()
    docx_path = _make_docx(tmp_path)

    word_to_pptx(str(docx_path), str(out))

    leftovers = list(out.glob("*.pdf"))
    assert leftovers == []


def test_word_to_pptx_respects_dpi(tmp_path):
    out = tmp_path / "out"
    out.mkdir()
    docx_path = _make_docx(tmp_path)

    output_path = word_to_pptx(str(docx_path), str(out), dpi=96)

    assert Path(output_path).exists()
    prs = Presentation(output_path)
    assert len(prs.slides) >= 1


# ──────────────────────────────────────────────────────────────
# Standalone API endpoint
# ──────────────────────────────────────────────────────────────

@pytest.fixture
def mock_dirs(tmp_path):
    upload_dir = tmp_path / "uploads"
    output_dir = tmp_path / "outputs"
    upload_dir.mkdir()
    output_dir.mkdir()
    with patch("main.UPLOAD_DIR", upload_dir), patch("main.OUTPUT_DIR", output_dir):
        yield {"upload": upload_dir, "output": output_dir}


def test_api_word_to_pptx(mock_dirs, auth_client, tmp_path):
    docx_path = _make_docx(tmp_path, name="Slides Source.docx")

    with open(docx_path, "rb") as f:
        files = {"file": (docx_path.name, f, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
        response = auth_client.post("/api/word/to-pptx", files=files)

    assert response.status_code == 200
    data = response.json()
    assert data["status"] == "success"
    assert data["filename"] == "Slides Source_forgefiles.org.pptx"
    assert result_path(mock_dirs["output"], data).exists()


def test_api_word_to_pptx_invalid_dpi(mock_dirs, auth_client, tmp_path):
    docx_path = _make_docx(tmp_path)

    with open(docx_path, "rb") as f:
        files = {"file": (docx_path.name, f, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
        response = auth_client.post("/api/word/to-pptx", files=files, data={"dpi": 5000})

    assert response.status_code == 422


# ──────────────────────────────────────────────────────────────
# Workflow builder step
# ──────────────────────────────────────────────────────────────

def test_api_workflow_word_to_pptx_step(mock_dirs, auth_client, tmp_path):
    docx_path = _make_docx(tmp_path)
    steps = json.dumps([{"type": "word_to_pptx", "config": {"dpi": 96}, "label": "Word to PowerPoint"}])

    with open(docx_path, "rb") as f:
        files = {"file": (docx_path.name, f, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
        response = auth_client.post("/api/workflow/execute", files=files, data={"steps": steps})

    assert response.status_code == 200
    assert "complete" in response.text
    assert "_forgefiles.org.pptx" in response.text
