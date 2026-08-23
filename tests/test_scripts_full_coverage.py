"""Branch-complete coverage for public/scripts helpers.

Covers the remaining edge branches: excel/ppt LibreOffice-fallback paths and
validation errors, image-util validation and rare-mode branches, the OCR engine
1.x/2.x API shims and factory double-check, event-log plumbing (get_connection,
writer-close failure, referrer parse failure), pdf_password_remover CLI error
paths, and seo_content/utils stragglers.
"""
import io
import sys
import threading
from types import SimpleNamespace

import pytest


# --------------------------------------------------------------------------
# scripts/utils.py
# --------------------------------------------------------------------------
class TestUtilsTryFont:
    def test_falls_back_to_default_font(self, monkeypatch):
        from PIL import ImageFont

        from scripts import utils

        # Pillow's load_default() itself calls truetype() on an embedded
        # font stream — only fail the named on-disk fonts.
        real_truetype = ImageFont.truetype

        def _no_font(name, size, **kwargs):
            if isinstance(name, str):
                raise OSError("no such font")
            return real_truetype(name, size, **kwargs)

        monkeypatch.setattr(ImageFont, "truetype", _no_font)
        font = utils.try_font(12)
        assert font is not None  # ImageFont.load_default()


# --------------------------------------------------------------------------
# scripts/excel_utils.py
# --------------------------------------------------------------------------
class TestExcelUtils:
    def test_sanitize_cell_none(self):
        from scripts import excel_utils
        assert excel_utils._sanitize_cell(None) == ""
        assert excel_utils._sanitize_cell(3.5) == "3.5"

    def test_libreoffice_output_renamed_to_branded(self, tmp_path, monkeypatch):
        from scripts import excel_utils

        def fake_lo(input_path, output_dir):
            produced = tmp_path / "lo_output.pdf"
            produced.write_bytes(b"%PDF-fake")
            return produced

        monkeypatch.setattr(excel_utils, "libreoffice_to_pdf", fake_lo)
        xlsx = tmp_path / "book.xlsx"
        xlsx.write_bytes(b"PK fake")
        out = excel_utils.excel_to_pdf(str(xlsx), str(tmp_path))
        assert out.endswith("book_forgefiles.org.pdf")
        assert not (tmp_path / "lo_output.pdf").exists()

    def test_workbook_without_sheets_raises(self, tmp_path, monkeypatch):
        from scripts import excel_utils

        class FakeWB:
            sheetnames = []

            def close(self):
                pass

        monkeypatch.setattr(excel_utils, "libreoffice_to_pdf", lambda *a: None)
        monkeypatch.setattr(excel_utils, "load_workbook", lambda *a, **k: FakeWB())
        xlsx = tmp_path / "empty.xlsx"
        xlsx.write_bytes(b"PK fake")
        with pytest.raises(ValueError, match="no sheets"):
            excel_utils.excel_to_pdf(str(xlsx), str(tmp_path))

    def test_reportlab_fallback_full_featured(self, tmp_path, monkeypatch):
        """Two sheets (PageBreak), an empty sheet, wide rows (truncation note),
        and None cells all in one fallback render."""
        from openpyxl import Workbook

        from scripts import excel_utils

        monkeypatch.setattr(excel_utils, "libreoffice_to_pdf", lambda *a: None)
        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws.append([None] + ["c%d" % i for i in range(25)])  # 26 cols > cap
        wb.create_sheet("Second").append(["x"])
        wb.create_sheet("Empty")  # no rows at all
        src = tmp_path / "full.xlsx"
        wb.save(src)
        wb.close()

        out = excel_utils.excel_to_pdf(str(src), str(tmp_path))
        assert out.endswith(".pdf")
        with open(out, "rb") as fh:
            assert fh.read(5) == b"%PDF-"

    def test_merge_requires_inputs(self):
        from scripts import excel_utils
        with pytest.raises(ValueError, match="No input files"):
            excel_utils.merge_excel_files([], "out")
        with pytest.raises(ValueError, match="at least two"):
            excel_utils.merge_excel_files(["one.xlsx"], "out")

    def test_merge_deduplicates_sheet_names(self, tmp_path):
        from openpyxl import Workbook

        from scripts import excel_utils

        paths = []
        for i in (1, 2):
            wb = Workbook()
            wb.active.title = "Sheet"
            wb.active.append([i])
            p = tmp_path / f"sub{i}" / "book.xlsx"  # same stem both times
            p.parent.mkdir()
            wb.save(p)
            wb.close()
            paths.append(str(p))
        out = excel_utils.merge_excel_files(paths, str(tmp_path))

        from openpyxl import load_workbook
        merged = load_workbook(out, read_only=True)
        assert len(merged.sheetnames) == 2
        assert len(set(merged.sheetnames)) == 2
        merged.close()

    def test_merge_all_empty_inputs_creates_placeholder(self, tmp_path, monkeypatch):
        from scripts import excel_utils

        class FakeWB:
            sheetnames = []

            def close(self):
                pass

        monkeypatch.setattr(excel_utils, "load_workbook", lambda *a, **k: FakeWB())
        out = excel_utils.merge_excel_files(["a.xlsx", "b.xlsx"], str(tmp_path))
        from openpyxl import load_workbook
        wb = load_workbook(out, read_only=True)
        assert wb.sheetnames == ["Sheet1"]
        wb.close()


# --------------------------------------------------------------------------
# scripts/ppt_utils.py
# --------------------------------------------------------------------------
class TestPptUtils:
    def test_emu_to_px_none(self):
        from scripts import ppt_utils
        assert ppt_utils._emu_to_px(None) == 0
        assert ppt_utils._emu_to_px(9525) == 1

    def _fake_slide(self, shapes):
        return SimpleNamespace(shapes=shapes)

    def test_shape_with_broken_position_renders_at_origin(self):
        from scripts import ppt_utils

        class BadShape:
            has_text_frame = False
            shape_type = 1

            @property
            def left(self):
                raise RuntimeError("no position")

        img = ppt_utils._render_slide_to_image(self._fake_slide([BadShape()]), 95250, 95250)
        assert img.size == (10, 10)

    def test_picture_shape_is_pasted(self):
        from PIL import Image

        from scripts import ppt_utils

        buf = io.BytesIO()
        Image.new("RGB", (4, 4), "red").save(buf, "PNG")

        pic_shape = SimpleNamespace(
            left=0, top=0, width=4 * 9525, height=4 * 9525,
            shape_type=13, has_text_frame=False,
            image=SimpleNamespace(blob=buf.getvalue()),
        )
        img = ppt_utils._render_slide_to_image(
            self._fake_slide([pic_shape]), 20 * 9525, 20 * 9525)
        assert img.getpixel((1, 1)) == (255, 0, 0)

    def test_picture_decode_failure_is_skipped(self):
        from scripts import ppt_utils

        pic_shape = SimpleNamespace(
            left=0, top=0, width=9525, height=9525,
            shape_type=13, has_text_frame=False,
            image=SimpleNamespace(blob=b"not an image"),
        )
        img = ppt_utils._render_slide_to_image(
            self._fake_slide([pic_shape]), 10 * 9525, 10 * 9525)
        assert img.getpixel((0, 0)) == (255, 255, 255)  # untouched background

    def test_empty_and_colored_paragraphs(self):
        from scripts import ppt_utils

        empty_para = SimpleNamespace(runs=[], text="")
        colored_run = SimpleNamespace(
            text="hello",
            font=SimpleNamespace(size=None, color=SimpleNamespace(rgb=(255, 0, 0))),
        )
        color_para = SimpleNamespace(runs=[colored_run], text="hello")
        text_shape = SimpleNamespace(
            left=0, top=0, width=100 * 9525, height=50 * 9525,
            shape_type=1, has_text_frame=True,
            text_frame=SimpleNamespace(paragraphs=[empty_para, color_para]),
        )
        img = ppt_utils._render_slide_to_image(
            self._fake_slide([text_shape]), 200 * 9525, 200 * 9525)
        assert img is not None  # both paragraph branches ran without error

    def test_libreoffice_output_renamed(self, tmp_path, monkeypatch):
        from scripts import ppt_utils

        def fake_lo(input_path, output_dir):
            produced = tmp_path / "lo_out.pdf"
            produced.write_bytes(b"%PDF-fake")
            return produced

        monkeypatch.setattr(ppt_utils, "libreoffice_to_pdf", fake_lo)
        pptx = tmp_path / "deck.pptx"
        pptx.write_bytes(b"PK fake")
        out = ppt_utils.ppt_to_pdf(str(pptx), str(tmp_path))
        assert out.endswith("deck_forgefiles.org.pdf")

    def test_presentation_without_slides_raises(self, tmp_path, monkeypatch):
        from pptx import Presentation

        from scripts import ppt_utils

        monkeypatch.setattr(ppt_utils, "libreoffice_to_pdf", lambda *a: None)
        src = tmp_path / "empty.pptx"
        Presentation().save(src)
        with pytest.raises(ValueError, match="no slides"):
            ppt_utils.ppt_to_pdf(str(src), str(tmp_path))

    def test_merge_requires_inputs(self):
        from scripts import ppt_utils
        with pytest.raises(ValueError, match="No input files"):
            ppt_utils.merge_pptx([], "out")
        with pytest.raises(ValueError, match="at least two"):
            ppt_utils.merge_pptx(["one.pptx"], "out")

    def test_merge_strips_layout_placeholders(self, tmp_path, monkeypatch):
        """Fully-faked Presentation: the layout adds a placeholder that must be
        stripped before the source slide's shape XML is copied in."""
        from scripts import ppt_utils

        class _El:
            def __init__(self):
                self.parent = None

            def getparent(self):
                return self.parent

        class _Parent:
            def __init__(self):
                self.removed = []

            def remove(self, el):
                self.removed.append(el)

        class _Shape:
            def __init__(self):
                self._element = _El()

        class _SpTree:
            def __init__(self):
                self.inserted = []

            def insert_element_before(self, el, marker):
                self.inserted.append(el)

        class _ShapeCollection(list):
            def __init__(self, *a):
                super().__init__(*a)
                self._spTree = _SpTree()

        class _Slide:
            def __init__(self, shapes=None):
                self.shapes = _ShapeCollection(shapes or [])

        class _Slides:
            def __init__(self):
                self._slides = []
                self.placeholder_parents = []

            def add_slide(self, layout):
                ph = _Shape()
                ph._element.parent = _Parent()
                self.placeholder_parents.append(ph._element.parent)
                slide = _Slide([ph])
                self._slides.append(slide)
                return slide

            def __iter__(self):
                return iter(self._slides)

        class FakePresentation:
            def __init__(self, path):
                self.slide_layouts = [object()] * 11
                self.slides = _Slides()
                if "src" in str(path):
                    self.slides._slides.append(_Slide([_Shape()]))

            def save(self, out):
                self.saved_to = out

        monkeypatch.setattr(ppt_utils, "Presentation", FakePresentation)
        out = ppt_utils.merge_pptx(
            [str(tmp_path / "base.pptx"), str(tmp_path / "src.pptx")], str(tmp_path))
        assert out.endswith(".pptx")


# --------------------------------------------------------------------------
# scripts/image_utils.py
# --------------------------------------------------------------------------
class TestImageUtils:
    def test_flatten_non_rgb_mode(self):
        from PIL import Image

        from scripts import image_utils
        out = image_utils._flatten_to_rgb(Image.new("L", (4, 4)))
        assert out.mode == "RGB"

    def test_preserved_save_kwargs_getexif_failure(self):
        from scripts import image_utils

        class NoExif:
            info = {}

            def getexif(self):
                raise RuntimeError("no exif support")

        assert image_utils._preserved_save_kwargs(NoExif()) == {}

    def test_preserved_save_kwargs_tobytes_failure(self):
        from scripts import image_utils

        class BadExif(dict):
            def tobytes(self):
                raise RuntimeError("broken exif")

        img = SimpleNamespace(info={}, getexif=lambda: BadExif({0x0112: 6}))
        assert image_utils._preserved_save_kwargs(img) == {}

    def _photo(self, tmp_path, size=(400, 300)):
        import random

        from PIL import Image
        rng = random.Random(7)
        img = Image.new("RGB", size)
        img.putdata([(rng.randrange(256), rng.randrange(256), rng.randrange(256))
                     for _ in range(size[0] * size[1])])
        p = tmp_path / "photo.jpg"
        img.save(p, "JPEG", quality=95)
        return p

    def test_resize_target_size_generous(self, tmp_path):
        from scripts import image_utils
        src = self._photo(tmp_path)
        out = image_utils.resize_image(str(src), str(tmp_path), "target_size",
                                       target_size_kb=5000)
        assert out.endswith(".jpg")

    def test_resize_target_size_tiny_hits_floor(self, tmp_path):
        from scripts import image_utils
        src = self._photo(tmp_path)
        out = image_utils.resize_image(str(src), str(tmp_path), "target_size",
                                       target_size_kb=1)
        assert out.endswith(".jpg")  # resized down to the 10px floor, then saved

    def test_rotate_unknown_extension_defaults_jpg(self, tmp_path):
        from PIL import Image

        from scripts import image_utils
        src = tmp_path / "pic.bmp"
        Image.new("RGB", (10, 10)).save(src)
        out = image_utils.rotate_image(str(src), str(tmp_path), 90)
        assert out.endswith(".jpg")

    def test_compress_quality_not_a_number(self, tmp_path):
        from scripts import image_utils
        with pytest.raises(ValueError, match="integer"):
            image_utils.compress_image("x.jpg", str(tmp_path), quality="high")

    def test_convert_png_palette_to_rgba(self, tmp_path):
        from PIL import Image

        from scripts import image_utils
        src = tmp_path / "pal.png"
        Image.new("P", (8, 8)).save(src)
        out = image_utils.convert_image_format(str(src), str(tmp_path), "png")
        assert out.endswith(".png")

    def test_watermark_opacity_not_a_number(self, tmp_path):
        from scripts import image_utils
        with pytest.raises(ValueError, match="opacity"):
            image_utils.watermark_image("x.jpg", str(tmp_path), "txt", opacity="lots")

    def test_watermark_textbbox_fallback(self, tmp_path, monkeypatch):
        from PIL import Image, ImageDraw, ImageFont

        from scripts import image_utils

        def _no_bbox(self, *a, **k):
            raise AttributeError("old pillow")

        monkeypatch.setattr(ImageDraw.ImageDraw, "textbbox", _no_bbox)
        # The legacy fallback calls font.getsize(), which current Pillow fonts
        # no longer have — wrap a real font so getsize exists again while
        # getmask & friends still delegate to it.
        class _OldFont:
            def __init__(self, font):
                self._font = font

            def getsize(self, text):
                return (10, 10)

            def __getattr__(self, name):
                return getattr(self._font, name)

        monkeypatch.setattr(
            image_utils, "try_font",
            lambda size: _OldFont(ImageFont.load_default()),
        )
        src = tmp_path / "photo.png"
        Image.new("RGB", (60, 60), "blue").save(src)
        out = image_utils.watermark_image(str(src), str(tmp_path), "demo")
        assert out.endswith(".png")


# --------------------------------------------------------------------------
# scripts/ocr_engine.py
# --------------------------------------------------------------------------
class TestOcrEngine:
    def test_to_plain_list_variants(self):
        from scripts import ocr_engine
        assert ocr_engine._to_plain_list(None) == []
        assert ocr_engine._to_plain_list([1, 2]) == [1, 2]
        assert ocr_engine._to_plain_list((3,)) == [3]

        class HasTolist:
            def tolist(self):
                return ["n"]

        assert ocr_engine._to_plain_list(HasTolist()) == ["n"]

    def test_rapidocr_legacy_package_fallback(self, monkeypatch):
        from scripts import ocr_engine

        class FakeRapidOCR:
            def __call__(self, img):
                # v1.x shape: (list of [box, text, score] detections, elapse)
                return ([[[[0, 0], [1, 0], [1, 1], [0, 1]], "hi", 0.9]], 0.01)

        monkeypatch.setitem(sys.modules, "rapidocr", None)  # force ImportError
        monkeypatch.setitem(sys.modules, "rapidocr_onnxruntime",
                            SimpleNamespace(RapidOCR=FakeRapidOCR))
        engine = ocr_engine.RapidOCREngine()
        out = engine.recognize("img.png")
        assert out[0]["text"] == "hi"
        assert out[0]["confidence"] == pytest.approx(0.9)
        assert engine.name == "rapidocr"

    def test_rapidocr_v2_output_skips_empty_text(self):
        from scripts import ocr_engine

        engine = ocr_engine.RapidOCREngine.__new__(ocr_engine.RapidOCREngine)
        engine._engine = lambda img: SimpleNamespace(
            txts=["", "real"], scores=[0.5], boxes=None)
        out = engine.recognize("img.png")
        assert out == [{"text": "real", "confidence": 0.0, "bbox": []}]

    def test_paddle_recognize_block_shapes(self):
        from scripts import ocr_engine

        engine = ocr_engine.PaddleOCREngine.__new__(ocr_engine.PaddleOCREngine)
        engine._engine = lambda img: [
            "not-a-dict",                                  # skipped
            {"res": [{"text": "cell", "confidence": 0.7,
                      "text_region": [[0, 0]]}]},          # layout block
            {"res": "not-a-list"},                         # ignored
            {"res": [{"text": ""}]},                       # empty text skipped
            {"text": "line", "score": 0.6, "bbox": [1, 2]},  # flat block
        ]
        out = engine.recognize("img.png")
        assert [i["text"] for i in out] == ["cell", "line"]
        assert out[1]["bbox"] == [1, 2]
        assert engine.name == "paddle"
        assert engine.supports_layout is True

    def test_factory_double_checked_locking(self, monkeypatch):
        from scripts import ocr_engine

        sentinel = object()

        class FlipLock:
            def __enter__(self):
                # Simulate another thread finishing initialization while we
                # waited on the lock.
                ocr_engine._ocr_engine = sentinel
                ocr_engine._ocr_engine_loaded = True
                return self

            def __exit__(self, *exc):
                return False

        monkeypatch.setattr(ocr_engine, "_ocr_engine", None)
        monkeypatch.setattr(ocr_engine, "_ocr_engine_loaded", False)
        monkeypatch.setattr(ocr_engine, "_ocr_engine_lock", FlipLock())
        assert ocr_engine.get_ocr_engine() is sentinel


# --------------------------------------------------------------------------
# scripts/event_log.py
# --------------------------------------------------------------------------
class TestEventLog:
    @pytest.fixture
    def event_db(self, tmp_path, monkeypatch):
        db = tmp_path / "events.db"
        monkeypatch.setenv("EVENT_DB_PATH", str(db))
        yield db
        from scripts import event_log
        event_log.close_connections()

    def test_get_connection_opens_readable_db(self, event_db):
        from scripts import event_log
        event_log.log_event("op", success=True, duration_ms=1)
        conn = event_log.get_connection()
        try:
            rows = conn.execute("SELECT COUNT(*) FROM operation_events").fetchone()
            assert rows[0] == 1
        finally:
            conn.close()

    def test_ensure_schema_double_check(self, tmp_path, monkeypatch):
        from scripts import event_log

        path = tmp_path / "fresh.db"
        event_log._initialized_paths.discard(str(path))

        real_lock = event_log._init_lock

        class FlipLock:
            def __enter__(self):
                event_log._initialized_paths.add(str(path))
                return self

            def __exit__(self, *exc):
                return False

        monkeypatch.setattr(event_log, "_init_lock", FlipLock())
        event_log._ensure_schema(path)  # returns via the in-lock re-check
        assert not path.exists()  # no schema was actually created
        monkeypatch.setattr(event_log, "_init_lock", real_lock)
        event_log._initialized_paths.discard(str(path))

    def test_migrate_columns_survives_operational_error(self):
        import sqlite3 as sql

        from scripts import event_log

        class FakeCursor:
            def __iter__(self):
                return iter([])

        class FakeConn:
            def execute(self, sql_text, *a):
                if sql_text.startswith("PRAGMA table_info"):
                    return FakeCursor()
                raise sql.OperationalError("readonly database")

        event_log._migrate_columns(FakeConn())  # must not raise

    def test_close_writer_survives_close_failure(self, monkeypatch):
        from scripts import event_log

        class BadConn:
            def close(self):
                raise RuntimeError("already closed")

        monkeypatch.setattr(event_log, "_writer", (None, BadConn()))
        event_log.close_connections()  # must not raise
        assert event_log._writer is None

    def test_sanitize_error_plain_value(self):
        from scripts import event_log
        assert event_log.sanitize_error("just a string") == "just a string"

    def test_referrer_host_unparseable_is_direct(self):
        from scripts import event_log
        assert event_log.referrer_host("http://[broken-ipv6") == event_log.REFERRER_DIRECT

    def test_unknown_funnel_event_ignored(self, event_db):
        from scripts import event_log
        assert event_log.log_funnel_event("not_a_real_event") is False


# --------------------------------------------------------------------------
# scripts/pdf_password_remover.py — CLI branches
# --------------------------------------------------------------------------
class TestPdfPasswordRemover:
    def _make_pdf(self, path, password=None):
        import pikepdf
        pdf = pikepdf.new()
        if password:
            pdf.save(path, encryption=pikepdf.Encryption(owner=password, user=password))
        else:
            pdf.save(path)
        pdf.close()

    def test_missing_pikepdf_reraises(self, tmp_path, monkeypatch):
        from scripts import pdf_password_remover as ppr
        src = tmp_path / "doc.pdf"
        self._make_pdf(src)
        monkeypatch.setitem(sys.modules, "pikepdf", None)
        with pytest.raises(ImportError):
            ppr.remove_pdf_password(str(src), "pw", str(tmp_path / "out.pdf"))

    def test_cli_empty_directory_exits_cleanly(self, tmp_path, monkeypatch, capsys):
        from scripts import pdf_password_remover as ppr
        monkeypatch.setattr(sys, "argv", ["prog", str(tmp_path), "pw"])
        with pytest.raises(SystemExit) as exc:
            ppr.main()
        assert exc.value.code == 0
        assert "No PDF files" in capsys.readouterr().out

    def test_cli_batch_wrong_password(self, tmp_path, monkeypatch, capsys):
        from scripts import pdf_password_remover as ppr
        self._make_pdf(tmp_path / "locked.pdf", password="right")
        monkeypatch.setattr(sys, "argv", ["prog", str(tmp_path), "wrong"])
        ppr.main()
        out = capsys.readouterr().out
        assert "Incorrect password" in out
        assert "Failed: 1" in out

    def test_cli_batch_password_error_without_pikepdf_importable(
        self, tmp_path, monkeypatch, capsys
    ):
        from scripts import pdf_password_remover as ppr
        (tmp_path / "f.pdf").write_bytes(b"%PDF-fake")

        PasswordError = type("PasswordError", (Exception,), {})

        def _raise(*a, **k):
            raise PasswordError("nope")

        monkeypatch.setattr(ppr, "remove_pdf_password", _raise)
        monkeypatch.setitem(sys.modules, "pikepdf", None)
        monkeypatch.setattr(sys, "argv", ["prog", str(tmp_path), "pw"])
        ppr.main()
        assert "Error: nope" in capsys.readouterr().out

    def test_cli_single_missing_pikepdf(self, tmp_path, monkeypatch, capsys):
        from scripts import pdf_password_remover as ppr
        src = tmp_path / "doc.pdf"
        self._make_pdf(src)

        def _raise(*a, **k):
            raise ImportError("No module named 'pikepdf'", name="pikepdf")

        monkeypatch.setattr(ppr, "remove_pdf_password", _raise)
        monkeypatch.setattr(sys, "argv", ["prog", str(src), "pw"])
        with pytest.raises(SystemExit):
            ppr.main()
        assert "pikepdf is not installed" in capsys.readouterr().out

    def test_cli_single_other_import_error_reraises(self, tmp_path, monkeypatch):
        from scripts import pdf_password_remover as ppr
        src = tmp_path / "doc.pdf"
        self._make_pdf(src)

        def _raise(*a, **k):
            raise ImportError("No module named 'other'", name="other")

        monkeypatch.setattr(ppr, "remove_pdf_password", _raise)
        monkeypatch.setattr(sys, "argv", ["prog", str(src), "pw"])
        with pytest.raises(ImportError):
            ppr.main()

    def test_cli_single_password_error_without_pikepdf_importable(
        self, tmp_path, monkeypatch, capsys
    ):
        from scripts import pdf_password_remover as ppr
        src = tmp_path / "doc.pdf"
        self._make_pdf(src)

        PasswordError = type("PasswordError", (Exception,), {})

        def _raise(*a, **k):
            raise PasswordError("bad pw")

        monkeypatch.setattr(ppr, "remove_pdf_password", _raise)
        monkeypatch.setitem(sys.modules, "pikepdf", None)
        monkeypatch.setattr(sys, "argv", ["prog", str(src), "pw"])
        with pytest.raises(SystemExit):
            ppr.main()
        assert "Error: bad pw" in capsys.readouterr().out


# --------------------------------------------------------------------------
# scripts/seo_content.py stragglers
# --------------------------------------------------------------------------
class TestSeoContent:
    def test_upload_box_plain_link_for_non_category_tool(self):
        from scripts import seo_content
        html = seo_content._upload_box("some-tool", {"tool": "unknown", "cta": "Go"})
        assert 'class="cta"' in html
        assert "data-ff-upload" not in html

    def test_all_tool_slugs(self):
        from scripts import seo_content
        slugs = seo_content.all_tool_slugs()
        assert slugs and "merge-pdf" in slugs
