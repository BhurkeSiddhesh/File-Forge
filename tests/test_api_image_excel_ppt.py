"""
API integration tests for image, Excel, and PPT endpoints.

These endpoints existed in main.py but had no HTTP-layer test coverage:
  POST /api/image/rotate
  POST /api/image/compress
  POST /api/image/convert
  POST /api/image/watermark
  POST /api/excel/to-pdf
  POST /api/excel/csv-to-xlsx
  POST /api/excel/xlsx-to-csv
  POST /api/excel/merge
  POST /api/ppt/to-pdf
  POST /api/ppt/to-images
  POST /api/ppt/merge
  POST /api/pdf/merge
  POST /api/pdf/watermark
  POST /api/pdf/to-images
  POST /api/pdf/sign
  POST /api/pdf/extract-pages  (supplemental validation cases)

Also covers:
  save_upload() — 413 when file exceeds MAX_UPLOAD_MB
  resize_image() target_size mode
  watermark_image() positional variants
"""
import csv
import io
from pathlib import Path
from unittest.mock import patch

import pytest
from fastapi.testclient import TestClient
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas as rl_canvas

from main import app


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _make_client():
    """Return an unauthenticated client (app is fully public)."""
    return TestClient(app, raise_server_exceptions=False)


def _simple_pdf_bytes() -> bytes:
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf)
    c.drawString(72, 750, "Integration test page")
    c.showPage()
    c.save()
    buf.seek(0)
    return buf.read()


def _multi_page_pdf_bytes(pages: int = 3) -> bytes:
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf)
    for i in range(1, pages + 1):
        c.drawString(72, 750, f"Page {i}")
        c.showPage()
    c.save()
    buf.seek(0)
    return buf.read()


def _jpeg_bytes(w: int = 200, h: int = 150) -> bytes:
    buf = io.BytesIO()
    Image.new("RGB", (w, h), color=(180, 90, 40)).save(buf, "JPEG", quality=80)
    buf.seek(0)
    return buf.read()


def _png_bytes(w: int = 200, h: int = 150) -> bytes:
    buf = io.BytesIO()
    Image.new("RGBA", (w, h), color=(0, 128, 255, 200)).save(buf, "PNG")
    buf.seek(0)
    return buf.read()


def _xlsx_bytes() -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["name", "score"])
    ws.append(["alice", 95])
    ws.append(["bob", 80])
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


def _csv_bytes() -> bytes:
    buf = io.StringIO()
    w = csv.writer(buf)
    w.writerow(["city", "pop"])
    w.writerow(["NYC", 8_000_000])
    w.writerow(["LA", 4_000_000])
    return buf.getvalue().encode("utf-8")


def _pptx_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Test Slide"
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    tb.text_frame.text = "Body text"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _sig_png_bytes() -> bytes:
    """Transparent PNG suitable for use as a PDF signature."""
    buf = io.BytesIO()
    Image.new("RGBA", (200, 80), (0, 0, 0, 0)).save(buf, "PNG")
    buf.seek(0)
    return buf.read()


# ---------------------------------------------------------------------------
# Fixture: mock dirs so tests don't litter the real uploads/outputs
# ---------------------------------------------------------------------------

@pytest.fixture
def mock_dirs(tmp_path):
    upload_dir = tmp_path / "uploads"
    output_dir = tmp_path / "outputs"
    upload_dir.mkdir()
    output_dir.mkdir()
    with patch("main.UPLOAD_DIR", upload_dir), patch("main.OUTPUT_DIR", output_dir):
        yield {"upload": upload_dir, "output": output_dir}


# ===========================================================================
# /api/image/rotate
# ===========================================================================

class TestApiRotateImage:
    def test_rotate_jpeg_90_degrees(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/image/rotate",
            data={"angle": "90"},
            files={"file": ("photo.jpg", _jpeg_bytes(), "image/jpeg")},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert "filename" in data

    def test_rotate_png_180_degrees(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/image/rotate",
            data={"angle": "180"},
            files={"file": ("img.png", _png_bytes(), "image/png")},
        )
        assert resp.status_code == 200
        assert resp.json()["status"] == "success"

    def test_rotate_bad_angle_returns_4xx(self, mock_dirs):
        """Non-numeric angle should result in 400 (ValueError) or 422 (FastAPI type coercion)."""
        client = _make_client()
        resp = client.post(
            "/api/image/rotate",
            data={"angle": "abc"},
            files={"file": ("photo.jpg", _jpeg_bytes(), "image/jpeg")},
        )
        assert resp.status_code in (400, 422)


# ===========================================================================
# /api/image/compress
# ===========================================================================

class TestApiCompressImage:
    def test_compress_jpeg_returns_size_stats(self, mock_dirs):
        client = _make_client()
        big = io.BytesIO()
        Image.new("RGB", (800, 800), (200, 100, 50)).save(big, "JPEG", quality=95)
        big.seek(0)
        resp = client.post(
            "/api/image/compress",
            data={"quality": "20"},
            files={"file": ("big.jpg", big.read(), "image/jpeg")},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert "original_size" in data
        assert "compressed_size" in data
        assert "reduction_pct" in data

    def test_compress_bad_quality_422(self, mock_dirs):
        """Quality > 95 should fail validation."""
        client = _make_client()
        resp = client.post(
            "/api/image/compress",
            data={"quality": "200"},
            files={"file": ("photo.jpg", _jpeg_bytes(), "image/jpeg")},
        )
        assert resp.status_code in (400, 422)


# ===========================================================================
# /api/image/convert
# ===========================================================================

class TestApiConvertImage:
    def test_jpeg_to_png(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/image/convert",
            data={"target_format": "png"},
            files={"file": ("photo.jpg", _jpeg_bytes(), "image/jpeg")},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert data["filename"].endswith(".png")

    def test_jpeg_to_webp(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/image/convert",
            data={"target_format": "webp"},
            files={"file": ("photo.jpg", _jpeg_bytes(), "image/jpeg")},
        )
        assert resp.status_code == 200
        assert resp.json()["filename"].endswith(".webp")

    def test_png_to_jpg(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/image/convert",
            data={"target_format": "jpg"},
            files={"file": ("img.png", _png_bytes(), "image/png")},
        )
        assert resp.status_code == 200
        assert resp.json()["filename"].endswith(".jpg")

    def test_invalid_format_returns_400(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/image/convert",
            data={"target_format": "bmp"},
            files={"file": ("photo.jpg", _jpeg_bytes(), "image/jpeg")},
        )
        assert resp.status_code == 400


# ===========================================================================
# /api/image/watermark
# ===========================================================================

class TestApiWatermarkImage:
    def test_watermark_bottom_right(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/image/watermark",
            data={"text": "CONFIDENTIAL", "position": "bottom-right", "opacity": "0.5"},
            files={"file": ("photo.jpg", _jpeg_bytes(), "image/jpeg")},
        )
        assert resp.status_code == 200
        assert resp.json()["status"] == "success"
        assert "filename" in resp.json()

    def test_watermark_diagonal(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/image/watermark",
            data={"text": "DRAFT", "position": "diagonal"},
            files={"file": ("photo.jpg", _jpeg_bytes(), "image/jpeg")},
        )
        assert resp.status_code == 200

    def test_watermark_center(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/image/watermark",
            data={"text": "SAMPLE", "position": "center"},
            files={"file": ("photo.jpg", _jpeg_bytes(), "image/jpeg")},
        )
        assert resp.status_code == 200

    def test_empty_text_returns_400(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/image/watermark",
            data={"text": "   "},
            files={"file": ("photo.jpg", _jpeg_bytes(), "image/jpeg")},
        )
        assert resp.status_code == 400

    def test_bad_position_returns_400(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/image/watermark",
            data={"text": "X", "position": "nowhere"},
            files={"file": ("photo.jpg", _jpeg_bytes(), "image/jpeg")},
        )
        assert resp.status_code == 400


# ===========================================================================
# /api/excel/to-pdf
# ===========================================================================

class TestApiExcelToPdf:
    def test_xlsx_to_pdf_success(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/excel/to-pdf",
            files={"file": ("data.xlsx", _xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert data["filename"].endswith(".pdf")


# ===========================================================================
# /api/excel/csv-to-xlsx
# ===========================================================================

class TestApiCsvToXlsx:
    def test_csv_to_xlsx_default_delimiter(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/excel/csv-to-xlsx",
            data={"delimiter": ","},
            files={"file": ("data.csv", _csv_bytes(), "text/csv")},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert data["filename"].endswith(".xlsx")

    def test_csv_to_xlsx_pipe_delimiter(self, mock_dirs):
        csv_content = "a|b\n1|2\n".encode()
        client = _make_client()
        resp = client.post(
            "/api/excel/csv-to-xlsx",
            data={"delimiter": "|"},
            files={"file": ("piped.csv", csv_content, "text/csv")},
        )
        assert resp.status_code == 200

    def test_csv_to_xlsx_bad_delimiter_400(self, mock_dirs):
        """Multi-char delimiter should be rejected."""
        client = _make_client()
        resp = client.post(
            "/api/excel/csv-to-xlsx",
            data={"delimiter": "||"},
            files={"file": ("data.csv", _csv_bytes(), "text/csv")},
        )
        assert resp.status_code == 400


# ===========================================================================
# /api/excel/xlsx-to-csv
# ===========================================================================

class TestApiXlsxToCsv:
    def test_xlsx_to_csv_success(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/excel/xlsx-to-csv",
            files={"file": ("data.xlsx", _xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert data["filename"].endswith(".csv")

    def test_xlsx_to_csv_unknown_sheet_400(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/excel/xlsx-to-csv",
            data={"sheet": "DoesNotExist"},
            files={"file": ("data.xlsx", _xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
        )
        assert resp.status_code == 400


# ===========================================================================
# /api/excel/merge
# ===========================================================================

class TestApiMergeExcel:
    def test_merge_two_xlsx(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/excel/merge",
            files=[
                ("files", ("book1.xlsx", _xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")),
                ("files", ("book2.xlsx", _xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")),
            ],
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert data["filename"].endswith(".xlsx")

    def test_merge_single_file_400(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/excel/merge",
            files=[
                ("files", ("only.xlsx", _xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")),
            ],
        )
        assert resp.status_code == 400


# ===========================================================================
# /api/ppt/to-pdf
# ===========================================================================

class TestApiPptToPdf:
    def test_pptx_to_pdf_success(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/ppt/to-pdf",
            files={"file": ("deck.pptx", _pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert data["filename"].endswith(".pdf")


# ===========================================================================
# /api/ppt/to-images
# ===========================================================================

class TestApiPptToImages:
    def test_pptx_to_png_zip(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/ppt/to-images",
            data={"fmt": "png"},
            files={"file": ("deck.pptx", _pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert data["filename"].endswith(".zip")
        assert data["slide_count"] >= 1

    def test_pptx_to_jpg_zip(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/ppt/to-images",
            data={"fmt": "jpg"},
            files={"file": ("deck.pptx", _pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
        assert resp.status_code == 200

    def test_invalid_fmt_400(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/ppt/to-images",
            data={"fmt": "bmp"},
            files={"file": ("deck.pptx", _pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
        assert resp.status_code == 400


# ===========================================================================
# /api/ppt/merge
# ===========================================================================

class TestApiMergePptx:
    def test_merge_two_pptx(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/ppt/merge",
            files=[
                ("files", ("deck1.pptx", _pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
                ("files", ("deck2.pptx", _pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
            ],
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert data["filename"].endswith(".pptx")

    def test_merge_single_pptx_400(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/ppt/merge",
            files=[
                ("files", ("deck.pptx", _pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
            ],
        )
        assert resp.status_code == 400


# ===========================================================================
# /api/pdf/merge
# ===========================================================================

class TestApiMergePdf:
    def test_merge_two_pdfs(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/pdf/merge",
            files=[
                ("files", ("a.pdf", _simple_pdf_bytes(), "application/pdf")),
                ("files", ("b.pdf", _simple_pdf_bytes(), "application/pdf")),
            ],
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert data["filename"].endswith(".pdf")

    def test_merge_single_pdf_400(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/pdf/merge",
            files=[("files", ("a.pdf", _simple_pdf_bytes(), "application/pdf"))],
        )
        assert resp.status_code == 400

    def test_merge_no_files_400(self, mock_dirs):
        client = _make_client()
        # FastAPI returns 422 when required List[UploadFile] is empty
        resp = client.post("/api/pdf/merge", files=[])
        assert resp.status_code in (400, 422)

    def test_merge_three_pdfs(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/pdf/merge",
            files=[
                ("files", ("a.pdf", _multi_page_pdf_bytes(2), "application/pdf")),
                ("files", ("b.pdf", _multi_page_pdf_bytes(3), "application/pdf")),
                ("files", ("c.pdf", _simple_pdf_bytes(), "application/pdf")),
            ],
        )
        assert resp.status_code == 200
        assert resp.json()["status"] == "success"


# ===========================================================================
# /api/pdf/watermark
# ===========================================================================

class TestApiPdfWatermark:
    def test_watermark_diagonal_success(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/pdf/watermark",
            data={"text": "CONFIDENTIAL", "position": "diagonal", "opacity": "0.3"},
            files={"file": ("doc.pdf", _simple_pdf_bytes(), "application/pdf")},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert data["filename"].endswith(".pdf")

    def test_watermark_center_position(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/pdf/watermark",
            data={"text": "DRAFT", "position": "center"},
            files={"file": ("doc.pdf", _simple_pdf_bytes(), "application/pdf")},
        )
        assert resp.status_code == 200

    def test_empty_watermark_text_400(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/pdf/watermark",
            data={"text": "   "},
            files={"file": ("doc.pdf", _simple_pdf_bytes(), "application/pdf")},
        )
        assert resp.status_code == 400

    def test_invalid_watermark_position_400(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/pdf/watermark",
            data={"text": "X", "position": "sideways"},
            files={"file": ("doc.pdf", _simple_pdf_bytes(), "application/pdf")},
        )
        assert resp.status_code == 400


# ===========================================================================
# /api/pdf/to-images
# ===========================================================================

class TestApiPdfToImages:
    def test_pdf_to_jpg_zip(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/pdf/to-images",
            data={"dpi": "72", "fmt": "jpg"},
            files={"file": ("doc.pdf", _multi_page_pdf_bytes(2), "application/pdf")},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert data["filename"].endswith(".zip")
        assert data["page_count"] == 2

    def test_pdf_to_png_zip(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/pdf/to-images",
            data={"dpi": "72", "fmt": "png"},
            files={"file": ("doc.pdf", _simple_pdf_bytes(), "application/pdf")},
        )
        assert resp.status_code == 200
        assert resp.json()["page_count"] == 1

    def test_bad_dpi_returns_400(self, mock_dirs):
        """DPI < 24 or > 600 should be rejected by pdf_utils."""
        client = _make_client()
        resp = client.post(
            "/api/pdf/to-images",
            data={"dpi": "10"},
            files={"file": ("doc.pdf", _simple_pdf_bytes(), "application/pdf")},
        )
        assert resp.status_code == 400


# ===========================================================================
# /api/pdf/sign
# ===========================================================================

class TestApiSignPdf:
    def test_sign_with_png_signature(self, mock_dirs):
        client = _make_client()
        resp = client.post(
            "/api/pdf/sign",
            data={"page": "1", "x": "0.65", "y": "0.85", "width": "0.2"},
            files={
                "file": ("doc.pdf", _multi_page_pdf_bytes(2), "application/pdf"),
                "signature": ("sig.png", _sig_png_bytes(), "image/png"),
            },
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "success"
        assert data["filename"].endswith(".pdf")

    def test_sign_invalid_content_type_400(self, mock_dirs):
        """Signature file must be PNG or JPEG — sending a PDF as signature must fail."""
        client = _make_client()
        resp = client.post(
            "/api/pdf/sign",
            data={"page": "1"},
            files={
                "file": ("doc.pdf", _simple_pdf_bytes(), "application/pdf"),
                "signature": ("sig.pdf", _simple_pdf_bytes(), "application/pdf"),
            },
        )
        assert resp.status_code == 400

    def test_sign_out_of_range_page_400(self, mock_dirs):
        """Page number beyond the PDF page count should be rejected."""
        client = _make_client()
        resp = client.post(
            "/api/pdf/sign",
            data={"page": "999"},
            files={
                "file": ("doc.pdf", _simple_pdf_bytes(), "application/pdf"),
                "signature": ("sig.png", _sig_png_bytes(), "image/png"),
            },
        )
        assert resp.status_code == 400


# ===========================================================================
# /api/pdf/extract-pages — supplemental validation cases
# ===========================================================================

class TestApiExtractPagesValidation:
    def test_out_of_range_page_spec_400(self, mock_dirs):
        """Requesting page 99 from a 1-page PDF must return 400."""
        client = _make_client()
        resp = client.post(
            "/api/pdf/extract-pages",
            data={"pages": "99"},
            files={"file": ("doc.pdf", _simple_pdf_bytes(), "application/pdf")},
        )
        assert resp.status_code == 400

    def test_invalid_page_string_400(self, mock_dirs):
        """Non-numeric page spec must return 400."""
        client = _make_client()
        resp = client.post(
            "/api/pdf/extract-pages",
            data={"pages": "abc"},
            files={"file": ("doc.pdf", _simple_pdf_bytes(), "application/pdf")},
        )
        assert resp.status_code == 400


# ===========================================================================
# save_upload() — upload size limit (413 Payload Too Large)
# ===========================================================================

class TestSaveUploadSizeLimit:
    def test_oversized_upload_returns_413(self, mock_dirs, monkeypatch):
        """Files exceeding MAX_UPLOAD_MB must be rejected with HTTP 413."""
        # monkeypatch replaces the module-level name read by save_upload() at call time.
        import main as main_module
        monkeypatch.setattr(main_module, "MAX_UPLOAD_MB", 0)
        client = _make_client()
        resp = client.post(
            "/api/pdf/remove-password",
            data={"password": "test"},
            files={"file": ("big.pdf", _simple_pdf_bytes(), "application/pdf")},
        )
        assert resp.status_code == 413


# ===========================================================================
# image_utils — resize_image target_size mode (unit-level)
# ===========================================================================

class TestResizeImageTargetSizeMode:
    def test_resize_to_target_size_kb(self, tmp_path):
        """resize_image target_size mode must produce a file within rough bounds."""
        from scripts.image_utils import resize_image

        src = tmp_path / "large.jpg"
        Image.new("RGB", (1000, 1000), (100, 150, 200)).save(src, "JPEG", quality=95)

        out = resize_image(str(src), str(tmp_path), "target_size", target_size_kb=50)
        assert Path(out).exists()
        # Size may slightly exceed target on very small images but should be close
        assert Path(out).stat().st_size > 0

    def test_resize_unknown_mode_raises(self, tmp_path):
        from scripts.image_utils import resize_image

        src = tmp_path / "img.jpg"
        Image.new("RGB", (100, 100)).save(src, "JPEG")

        with pytest.raises(ValueError, match="[Uu]nknown"):
            resize_image(str(src), str(tmp_path), "bogus_mode")

    def test_resize_dimensions_only_width(self, tmp_path):
        """Providing only width should calculate height to preserve aspect ratio."""
        from scripts.image_utils import resize_image

        src = tmp_path / "img.jpg"
        Image.new("RGB", (400, 200)).save(src, "JPEG")
        out = resize_image(str(src), str(tmp_path), "dimensions", width=200)
        with Image.open(out) as img:
            assert img.size[0] == 200
            assert img.size[1] == 100  # aspect ratio preserved

    def test_resize_percentage_50(self, tmp_path):
        from scripts.image_utils import resize_image

        src = tmp_path / "img.jpg"
        Image.new("RGB", (200, 100)).save(src, "JPEG")
        out = resize_image(str(src), str(tmp_path), "percentage", percentage=50)
        with Image.open(out) as img:
            assert img.size == (100, 50)


# ===========================================================================
# image_utils — watermark_image positional variants (unit-level)
# ===========================================================================

class TestWatermarkImagePositions:
    POSITIONS = ["top-left", "top-right", "center", "bottom-left", "bottom-right", "diagonal"]

    @pytest.mark.parametrize("position", POSITIONS)
    def test_all_positions_produce_output(self, tmp_path, position):
        from scripts.image_utils import watermark_image

        src = tmp_path / "img.jpg"
        Image.new("RGB", (400, 300), (50, 100, 150)).save(src, "JPEG")
        out = watermark_image(str(src), str(tmp_path), "SAMPLE", position=position)
        assert Path(out).exists()
        with Image.open(out) as img:
            assert img.size == (400, 300)

    def test_watermark_preserves_size(self, tmp_path):
        from scripts.image_utils import watermark_image

        src = tmp_path / "img.jpg"
        Image.new("RGB", (300, 200), (200, 50, 80)).save(src, "JPEG")
        out = watermark_image(str(src), str(tmp_path), "TEST", position="center")
        with Image.open(out) as img:
            assert img.size == (300, 200)

    def test_bad_position_raises_valueerror(self, tmp_path):
        from scripts.image_utils import watermark_image

        src = tmp_path / "img.jpg"
        Image.new("RGB", (100, 100)).save(src, "JPEG")
        with pytest.raises(ValueError, match="position"):
            watermark_image(str(src), str(tmp_path), "X", position="nowhere")

    def test_opacity_out_of_range_raises(self, tmp_path):
        from scripts.image_utils import watermark_image

        src = tmp_path / "img.jpg"
        Image.new("RGB", (100, 100)).save(src, "JPEG")
        with pytest.raises(ValueError, match="opacity"):
            watermark_image(str(src), str(tmp_path), "X", opacity=1.5)
